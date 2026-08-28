import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import io
import os
import textwrap
import hashlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

st.set_page_config(page_title="Unified 9-Box Margin Analysis", layout="wide")


# ==========================================
# DATA LOADING FUNCTIONS
# ==========================================
def load_9box_data(file_path):
    try:
        xls = pd.ExcelFile(file_path)
        df_local = pd.read_excel(xls, sheet_name='LOCAL')
        df_export = pd.read_excel(xls, sheet_name='EXPORT')

        df_local['Source_Sheet'] = 'LOCAL'
        df_export['Source_Sheet'] = 'EXPORT'

        df_combined = pd.concat([df_local, df_export], ignore_index=True)

        numeric_cols = [
            'Gross Sales (Current)', 'Return (Current)', 'COGS_Regular (Current)', 'Royalty (Current)',
            'Gross Profit (Current)', 'Contribution Margin (Current)',
            'Qty (Current)', 'Gross Sales (Previous)', 'Qty (Previous)',
            'Amount FG', 'Amount Material'
        ]
        for col in numeric_cols:
            if col in df_combined.columns:
                df_combined[col] = pd.to_numeric(df_combined[col], errors='coerce').fillna(0)

        if 'Remark2' not in df_combined.columns:
            df_combined['Remark2'] = ''

        # Ensure key columns are read as strings to prevent scientific notation formatting
        text_cols = ['SKU', 'Produk_Key', 'Product Name', 'New Code', 'New Product Name', 'Remark', 'Remark2', 'Status',
                     'Country']
        for col in text_cols:
            if col in df_combined.columns:
                df_combined[col] = df_combined[col].astype(str).str.replace(r'\.0$', '',
                                                                            regex=True).str.strip().replace(
                    ['nan', 'NaN', 'None'], '')

        # --- GLOBAL FALLBACK MAPPING ---
        # Prevents empty SKUs from clumping into a single Pseudo-SKU during consolidation
        if 'New Code' in df_combined.columns:
            fallback_code = df_combined['Produk_Key'] if 'Produk_Key' in df_combined.columns else df_combined['SKU']
            df_combined['New Code'] = df_combined['New Code'].replace([r'^\s*$', 'nan', 'NaN', 'None', 'UNKNOWN'],
                                                                      np.nan, regex=True)
            df_combined['New Code'] = df_combined['New Code'].fillna(fallback_code).astype(str)

        if 'New Product Name' in df_combined.columns:
            fallback_name = df_combined['Product Name'] if 'Product Name' in df_combined.columns else df_combined['SKU']
            df_combined['New Product Name'] = df_combined['New Product Name'].replace(
                [r'^\s*$', 'nan', 'NaN', 'None', 'UNKNOWN'], np.nan, regex=True)
            df_combined['New Product Name'] = df_combined['New Product Name'].fillna(fallback_name).astype(str)

        # GROSS MARGIN FORMULA (Based on agreed pure mapping rule)
        if all(col in df_combined.columns for col in
               ['Gross Sales (Current)', 'Return (Current)', 'COGS_Regular (Current)', 'Royalty (Current)']):
            df_combined['Gross Margin (Current)'] = (
                    df_combined['Gross Sales (Current)'] +
                    df_combined['Return (Current)'] -
                    df_combined['COGS_Regular (Current)'] -
                    df_combined['Royalty (Current)']
            )

            df_combined['Gross Margin (%)'] = np.where(
                (df_combined['Gross Sales (Current)'] + df_combined['Return (Current)']) != 0,
                (df_combined['Gross Margin (Current)'] / (
                        df_combined['Gross Sales (Current)'] + df_combined['Return (Current)'])) * 100,
                0.0
            )

        if 'Qty Growth (%)' not in df_combined.columns:
            if 'Qty (Current)' in df_combined.columns and 'Qty (Previous)' in df_combined.columns:
                df_combined['Qty Growth (%)'] = np.where(
                    df_combined['Qty (Previous)'] != 0,
                    ((df_combined['Qty (Current)'] - df_combined['Qty (Previous)']) / df_combined[
                        'Qty (Previous)'].abs()) * 100,
                    0.0
                )
            else:
                df_combined['Qty Growth (%)'] = 0.0

        if 'Gross Profit (%)' not in df_combined.columns and 'Gross Profit (Current)' in df_combined.columns and 'Gross Sales (Current)' in df_combined.columns:
            df_combined['Gross Profit (%)'] = np.where(
                df_combined['Gross Sales (Current)'] > 0,
                (df_combined['Gross Profit (Current)'] / df_combined['Gross Sales (Current)']) * 100,
                np.where(df_combined['Gross Profit (Current)'] < 0, -100.0, 0.0)
            )

        if 'Contribution Margin (%)' not in df_combined.columns and 'Contribution Margin (Current)' in df_combined.columns and 'Gross Sales (Current)' in df_combined.columns:
            df_combined['Contribution Margin (%)'] = np.where(
                df_combined['Gross Sales (Current)'] > 0,
                (df_combined['Contribution Margin (Current)'] / df_combined['Gross Sales (Current)']) * 100,
                np.where(df_combined['Contribution Margin (Current)'] < 0, -100.0, 0.0)
            )

        return df_combined
    except Exception as e:
        st.error(f"Failed to load data: {e}")
        return pd.DataFrame()


# ==========================================
# 9-BOX SUMMARY GRID RENDERER
# ==========================================
def render_9box_summary_grid(df, margin_type, margin_val_col, y_low_thresh, y_high_thresh, x_axis_metric, x_low_thresh,
                             x_high_thresh):
    abbr = "GP" if margin_type == "Gross Profit" else ("GM" if margin_type == "Gross Margin" else "CM")

    def get_box_html(b_name, bg_color):
        if 'Dynamic 9-Box Category' not in df.columns:
            return f'<div class="box {bg_color}"><div class="box-title">{b_name}</div><div class="box-main">0</div></div>'

        sub = df[df['Dynamic 9-Box Category'].str.startswith(b_name)]
        cnt = len(sub)
        sales = sub['Gross Sales (Current)'].sum() if 'Gross Sales (Current)' in sub.columns else 0
        margin = sub[margin_val_col].sum() if margin_val_col in sub.columns else 0
        qty = sub['Qty (Current)'].sum() if 'Qty (Current)' in sub.columns else 0

        pct = (margin / sales * 100) if sales > 0 else 0
        margin_bn = margin / 1e9
        sales_bn = sales / 1e9

        return f'<div class="box {bg_color}"><div class="box-title">{b_name}</div><div class="box-main">{cnt:,} <span class="box-sub">SKUs</span></div><div class="box-sub-small">IDR {margin_bn:,.0f} Bn margin &middot; {pct:.0f}% {abbr}<br>Sales: IDR {sales_bn:,.0f} Bn &middot; Qty: {qty:,.0f}</div></div>'

    def format_x_val(val):
        if x_axis_metric == "Qty Growth (%)":
            return f"{val:.1f}%"
        else:
            if val >= 1e9 or val <= -1e9:
                return f"Rp{val / 1e9:,.1f}M"
            else:
                return f"Rp{val / 1e6:,.0f}Jt"

    x_hdr = "SALES GROWTH" if x_axis_metric == "Qty Growth (%)" else "GROSS SALES (ABSOLUTE)"

    html_content = textwrap.dedent(f"""
        <style>
        .grid-container {{ display: grid; grid-template-columns: 100px 1fr 1fr 1fr; grid-template-rows: 25px 25px 1fr 1fr 1fr; gap: 12px; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin-bottom: 2rem; margin-top: 1rem; }}
        .box {{ padding: 15px; color: white; border-radius: 2px; display: flex; flex-direction: column; justify-content: space-between; min-height: 145px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }}
        .box-title {{ font-size: 13px; font-weight: 600; opacity: 0.95; }}
        .box-main {{ font-size: 32px; font-weight: 700; margin: 8px 0; }}
        .box-sub {{ font-size: 13px; font-weight: 400; opacity: 0.9; }}
        .box-sub-small {{ font-size: 12px; font-weight: 500; opacity: 0.95; line-height: 1.5; }}
        .bg-blue {{ background-color: #3871b6; }}
        .bg-green {{ background-color: #319b5e; }}
        .bg-red {{ background-color: #c03d32; }}
        .bg-gold {{ background-color: #d89f0e; }}
        .hdr-main {{ grid-column: 1 / 5; text-align: center; font-weight: 700; font-size: 13px; color: #4b5563; letter-spacing: 1px; }}
        .hdr-col {{ text-align: center; color: #6b7280; font-size: 13px; font-weight: 700; align-self: end; }}
        .hdr-row {{ text-align: center; font-size: 12px; font-weight: 700; color: #6b7280; display: flex; flex-direction: column; justify-content: center; }}
        .axis-label-y {{ writing-mode: vertical-rl; transform: rotate(180deg); position: absolute; left: 0; font-size: 12px; letter-spacing: 1px; color: #4b5563; font-weight: 700; height: 100%; text-align: center; }}
        .legend-container {{ display: flex; justify-content: center; gap: 40px; margin-top: 10px; font-size: 13px; font-weight: 700; color: #374151; }}
        .legend-item {{ display: flex; align-items: center; gap: 8px; }}
        .legend-box {{ width: 25px; height: 15px; border-radius: 2px; }}
        </style>

        <div style="position: relative;">
        <div class="axis-label-y">{margin_type.upper()}</div>
        <div class="grid-container" style="margin-left: 30px;">
        <div class="hdr-main">{x_hdr} &rarr;</div>
        <div></div>
        <div class="hdr-col">Low<br>(&lt;{format_x_val(x_low_thresh)})</div>
        <div class="hdr-col">Med<br>({format_x_val(x_low_thresh)} &ndash; {format_x_val(x_high_thresh)})</div>
        <div class="hdr-col">High<br>(&gt;{format_x_val(x_high_thresh)})</div>

        <div class="hdr-row">High<br>&gt;{y_high_thresh:.1f}%</div>
        {get_box_html('Box 1', 'bg-blue')}
        {get_box_html('Box 4', 'bg-blue')}
        {get_box_html('Box 7', 'bg-green')}

        <div class="hdr-row">Med<br>{y_low_thresh:.1f}&ndash;{y_high_thresh:.1f}%</div>
        {get_box_html('Box 2', 'bg-gold')}
        {get_box_html('Box 5', 'bg-gold')}
        {get_box_html('Box 8', 'bg-green')}

        <div class="hdr-row">Low<br>&lt;{y_low_thresh:.1f}%</div>
        {get_box_html('Box 3', 'bg-red')}
        {get_box_html('Box 6', 'bg-gold')}
        {get_box_html('Box 9', 'bg-gold')}
        </div>
        </div>

        <div class="legend-container">
        <div class="legend-item"><div class="legend-box bg-green"></div> Grow</div>
        <div class="legend-item"><div class="legend-box bg-blue"></div> Keep</div>
        <div class="legend-item"><div class="legend-box bg-gold"></div> Fix</div>
        <div class="legend-item"><div class="legend-box bg-red"></div> Exit candidate</div>
        </div>
        <br>
    """)
    st.markdown(html_content, unsafe_allow_html=True)
    return html_content


# ==========================================
# PPTX GENERATION ENGINE
# ==========================================
def create_portfolio_presentation(df_main, margin_type, margin_val_col):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(22, 54, 92)
    GRAY_TEXT = RGBColor(89, 89, 89)
    ACCENT_BLUE = RGBColor(0, 112, 192)
    LIGHT_BG = RGBColor(242, 242, 242)

    def add_consulting_slide(title_text, subtitle_text, metrics, insights):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_TEXT

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = NAVY
        line.line.color.rgb = NAVY

        box_width = 4.0
        for i, metric in enumerate(metrics):
            left_pos = 0.5 + (i * (box_width + 0.165))
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.3),
                                           Inches(box_width), Inches(0.6))
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_BG
            shape.line.color.rgb = ACCENT_BLUE

            tf_m = shape.text_frame
            tf_m.vertical_anchor = 3
            p_m = tf_m.paragraphs[0]
            p_m.alignment = PP_ALIGN.CENTER
            p_m.text = metric
            p_m.font.bold = True
            p_m.font.size = Pt(14)
            p_m.font.color.rgb = NAVY

        chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(7.5), Inches(4.8))
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = LIGHT_BG
        chart_box.line.color.rgb = GRAY_TEXT

        tf_c = chart_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = f"[ Insert {margin_type} Matrix Visualization Here ]\n(Copy-paste directly from the dashboard)"
        p_c.font.size = Pt(14)
        p_c.font.color.rgb = GRAY_TEXT
        p_c.font.italic = True

        insight_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.6), Inches(5.0))
        tf_insight = insight_box.text_frame
        tf_insight.word_wrap = True

        p_ins_head = tf_insight.paragraphs[0]
        p_ins_head.text = "Executive Insights & Action Mandate"
        p_ins_head.font.bold = True
        p_ins_head.font.size = Pt(16)
        p_ins_head.font.color.rgb = NAVY

        for insight in insights:
            parts = insight.split(":", 1)
            p_i = tf_insight.add_paragraph()
            p_i.space_before = Pt(12)
            p_i.level = 0

            if len(parts) == 2:
                run_bold = p_i.add_run()
                run_bold.text = f"• {parts[0]}:"
                run_bold.font.bold = True
                run_bold.font.size = Pt(12)
                run_bold.font.color.rgb = NAVY

                run_text = p_i.add_run()
                run_text.text = parts[1]
                run_text.font.size = Pt(12)
                run_text.font.color.rgb = GRAY_TEXT
            else:
                run_text = p_i.add_run()
                run_text.text = f"• {insight}"
                run_text.font.size = Pt(12)
                run_text.font.color.rgb = GRAY_TEXT

        return slide

    total_rev = df_main['Gross Sales (Current)'].sum()
    total_margin = df_main[margin_val_col].sum()
    blended_margin = (total_margin / total_rev * 100) if total_rev > 0 else 0
    total_skus = len(df_main)

    rev_str = f"Rp {total_rev / 1e12:.2f} Trillion" if total_rev >= 1e12 else f"Rp {total_rev / 1e9:.2f} Billion"

    add_consulting_slide(
        title_text=f"Global Portfolio Rationalization ({margin_type} View)",
        subtitle_text="Comprehensive assessment highlighting significant margin dilution and strategic resource allocation.",
        metrics=[
            f"Total Active SKUs: {total_skus:,}",
            f"Total Revenue: {rev_str}",
            f"Blended {margin_type}: {blended_margin:.1f}%"
        ],
        insights=[
            "Value-Destructive Tail: A significant portion of SKUs exhibit severe margin compression. Immediate rationalization or price restructuring is mandated.",
            "Aggressive Pruning & Pricing Gate: Execute targeted price hikes for underperformers. SKUs failing to clear the profitability hurdle post-hike must be slated for phase-out.",
            "Core Growth Engine: Flagship SKUs continue to deliver superior margins. Capital allocation should immediately pivot toward expanding distribution for these core assets."
        ]
    )

    return prs


# ==========================================
# GLOBAL HELPER FUNCTION FOR EXPORT
# ==========================================
def get_raw_excel(df_in):
    df_out = df_in.drop(
        columns=['Bubble_Size', 'X_Plot', 'Y_Plot', 'Plot_Color_Category', 'Dynamic 9-Box Category', 'Venn_Box'],
        errors='ignore').copy()

    text_cols_prefixes = ['SKU', 'Produk_Key', 'Product Name', 'New Code', 'New Product Name', 'Remark', 'Remark2',
                          'Status', 'Country', 'Source_Sheet']
    for col in df_out.columns:
        if any(col.startswith(prefix) for prefix in text_cols_prefixes):
            df_out[col] = df_out[col].astype(str).replace({'nan': '', 'None': ''})

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        df_out.to_excel(writer, index=False, sheet_name='Export_Data')
        worksheet = writer.sheets['Export_Data']
        for col in worksheet.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = (max_length + 2)
            worksheet.column_dimensions[column].width = adjusted_width

    return buffer.getvalue()


def get_pill_path(x0, y0, x1, y1):
    if x1 - x0 > y1 - y0:
        r = (y1 - y0) / 2
        return f"M {x0 + r},{y0} L {x1 - r},{y0} A {r},{r} 0 0,1 {x1 - r},{y1} L {x0 + r},{y1} A {r},{r} 0 0,1 {x0 + r},{y0} Z"
    else:
        r = (x1 - x0) / 2
        return f"M {x0},{y0 + r} L {x0},{y1 - r} A {r},{r} 0 0,1 {x1},{y1 - r} L {x1},{y0 + r} A {r},{r} 0 0,1 {x0},{y0 + r} Z"


# ==========================================
# MAIN APPLICATION LOGIC
# ==========================================
def main():
    st.markdown("""
        <div style="background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%); 
                    padding: 1.2rem; border-radius: 15px; margin-bottom: 1rem; color: white;
                    text-align: center;">
            <h1 style="margin:0; color: #38bdf8; font-size: 1.8rem; font-weight: 800;">
                📊 UNIFIED 9-BOX PORTFOLIO MATRIX
            </h1>
            <p style="margin: 0.4rem 0 0 0; opacity: 0.9; font-size: 1rem;">
                Strategic Engine for Gross Margin & Contribution Margin Optimization
            </p>
        </div>
    """, unsafe_allow_html=True)

    # ---------------------------------------------------------
    # FILE UPLOADER
    # ---------------------------------------------------------
    st.markdown("### 📂 Upload Source Data")
    st.info(
        "Please upload the **9-Box Data Excel file** (e.g., Data 9-box CM.xlsx) from your local computer to initiate the analysis.")

    uploaded_file = st.file_uploader("Upload 9-Box Data File (.xlsx)", type=["xlsx"])

    if uploaded_file is None:
        st.stop()

    df_raw = load_9box_data(uploaded_file)

    if df_raw.empty:
        st.warning(
            "⚠️ File cannot be read or has an incorrect format. Ensure the file contains 'LOCAL' and 'EXPORT' sheets.")
        st.stop()

    df_main = df_raw.copy()

    # ---------------------------------------------------------
    # SKU CONSOLIDATION VIEW TOGGLE
    # ---------------------------------------------------------
    st.markdown("### 🔀 SKU Consolidation View")
    view_mode = st.radio(
        "Select Data Granularity:",
        [
            "Uncommon (Original Granular SKUs)",
            "Common (Consolidated Master SKUs)",
            "Commonized Only (Master C- Prefix SKUs)",
            "Uncommonized Only (Granular C- Prefix SKUs)"
        ],
        horizontal=True,
        help="Switch to 'Common' to group data by Master SKU. Choose 'Commonized Only' for merged items, or 'Uncommonized Only' to see the original granular details of those merged items."
    )

    if view_mode in ["Common (Consolidated Master SKUs)", "Commonized Only (Master C- Prefix SKUs)"]:
        if 'New Code' in df_main.columns and 'New Product Name' in df_main.columns:
            group_cols = ['Source_Sheet', 'New Code', 'New Product Name']

            num_cols = df_main.select_dtypes(include=[np.number]).columns.tolist()
            sum_cols = [c for c in num_cols if not c.endswith('(%)')]

            agg_dict = {c: 'sum' for c in sum_cols}

            all_cols = df_main.columns.tolist()

            special_cols = [c for c in ['Produk_Key', 'Product Name'] if c in all_cols]
            cat_cols = [c for c in all_cols if c not in num_cols and c not in group_cols and c not in special_cols]

            for cat in cat_cols:
                agg_dict[cat] = lambda x: next(iter(pd.Series(x).dropna()), '')

            for sp in special_cols:
                agg_dict[sp] = lambda x: list(pd.Series(x).dropna().astype(str).unique())

            df_main = df_main.groupby(group_cols, as_index=False).agg(agg_dict)

            expanded_col_names = []
            for sp in special_cols:
                if sp in df_main.columns:
                    exp_df = pd.DataFrame(df_main[sp].tolist(), index=df_main.index).fillna('')
                    exp_df.columns = [f"{sp}_{i + 1}" for i in range(exp_df.shape[1])]
                    df_main = pd.concat([df_main.drop(columns=[sp]), exp_df], axis=1)
                    expanded_col_names.extend(exp_df.columns.tolist())

            final_cols = [c for c in df_main.columns if c not in expanded_col_names] + expanded_col_names
            df_main = df_main[final_cols]

            if 'Qty (Current)' in df_main.columns and 'Qty (Previous)' in df_main.columns:
                df_main['Qty Growth (%)'] = np.where(
                    df_main['Qty (Previous)'] != 0,
                    ((df_main['Qty (Current)'] - df_main['Qty (Previous)']) / df_main['Qty (Previous)'].abs()) * 100,
                    0.0
                )

            if 'Gross Sales (Current)' in df_main.columns:
                sales = df_main['Gross Sales (Current)']
                if 'Gross Profit (Current)' in df_main.columns:
                    df_main['Gross Profit (%)'] = np.where(sales > 0, (df_main['Gross Profit (Current)'] / sales) * 100,
                                                           np.where(df_main['Gross Profit (Current)'] < 0, -100.0, 0.0))
                if 'Gross Margin (Current)' in df_main.columns:
                    df_main['Gross Margin (%)'] = np.where(sales > 0, (df_main['Gross Margin (Current)'] / sales) * 100,
                                                           np.where(df_main['Gross Margin (Current)'] < 0, -100.0, 0.0))
                if 'Contribution Margin (Current)' in df_main.columns:
                    df_main['Contribution Margin (%)'] = np.where(sales > 0, (
                            df_main['Contribution Margin (Current)'] / sales) * 100,
                                                                  np.where(df_main['Contribution Margin (Current)'] < 0,
                                                                           -100.0, 0.0))

        else:
            st.warning("⚠️ Columns 'New Code' and 'New Product Name' are missing in the uploaded dataset.")

    if view_mode in ["Common (Consolidated Master SKUs)", "Commonized Only (Master C- Prefix SKUs)"]:
        global_name_col = 'New Product Name' if 'New Product Name' in df_main.columns else 'Product Name'
    else:
        global_name_col = 'Product Name' if 'Product Name' in df_main.columns else (
            'SKU' if 'SKU' in df_main.columns else df_main.columns[0])

    st.markdown("---")

    tab_main_matrix, tab_b3_intersection, tab_progress = st.tabs([
        "📈 MAIN 9-BOX ANALYSIS",
        "📊 VENN DIAGRAM INFOGRAPHIC",
        "📉 RATIONALIZATION PROGRESS"
    ])

    # ---------------------------------------------------------
    # TAB 1: MAIN 9-BOX ANALYSIS
    # ---------------------------------------------------------
    with tab_main_matrix:
        st.markdown("### 🔍 Filter & Metric Configuration")

        col1, col2, col3 = st.columns(3)

        with col1:
            market_filter = st.selectbox(
                "Select Market Focus:",
                ["ALL (Local + Export)", "LOCAL", "EXPORT"]
            )

        with col2:
            margin_selector = st.selectbox(
                "Y-Axis Metric (Margin):",
                ["Gross Profit", "Gross Margin", "Contribution Margin"]
            )

        with col3:
            x_axis_selector = st.selectbox(
                "X-Axis Metric (Performance):",
                ["Gross Sales (Current)", "Qty Growth (%)"],
                help="Select 'Gross Sales (Current)' for absolute nominals, or 'Qty Growth (%)' to view growth rates."
            )

        col4, col5, col6, col7 = st.columns(4)

        with col4:
            bubble_size_selector = st.selectbox(
                "Bubble Size Metric:",
                ["Gross Sales (Current)", "Gross Profit (Current)", "Gross Margin (Current)",
                 "Contribution Margin (Current)", "Qty (Current)"],
                index=4
            )

        with col5:
            if 'Remark' in df_main.columns:
                valid_remarks = sorted(
                    [str(x) for x in df_main['Remark'].unique() if pd.notna(x) and str(x).strip() != ''])
                remark_options = ["ALL"] + valid_remarks
            else:
                remark_options = ["ALL"]

            remark_filter = st.multiselect("Select Remark(s):", remark_options, default=["ALL"])

        with col6:
            if 'Status' in df_main.columns:
                valid_statuses = sorted(
                    [str(x) for x in df_main['Status'].unique() if pd.notna(x) and str(x).strip() != ''])
                status_options = ["ALL"] + valid_statuses
            else:
                status_options = ["ALL"]

            status_filter = st.multiselect("Select Status(es):", status_options, default=["ALL"])

        with col7:
            sku_search = st.text_input("Specific SKU Search:", "")

        if margin_selector == "Gross Profit":
            y_col = 'Gross Profit (%)'
            margin_val_col = 'Gross Profit (Current)'
        elif margin_selector == "Gross Margin":
            y_col = 'Gross Margin (%)'
            margin_val_col = 'Gross Margin (Current)'
        else:
            y_col = 'Contribution Margin (%)'
            margin_val_col = 'Contribution Margin (Current)'

        x_col = x_axis_selector

        filtered_df = df_main.copy()

        if market_filter != "ALL (Local + Export)":
            filtered_df = filtered_df[filtered_df['Source_Sheet'] == market_filter]

        if remark_filter and 'Remark' in filtered_df.columns:
            actual_remarks = [r for r in remark_filter if r != "ALL"]
            if actual_remarks:
                filtered_df = filtered_df[filtered_df['Remark'].astype(str).isin(actual_remarks)]

        if status_filter and 'Status' in filtered_df.columns:
            actual_statuses = [s for s in status_filter if s != "ALL"]
            if actual_statuses:
                filtered_df = filtered_df[filtered_df['Status'].astype(str).isin(actual_statuses)]

        if sku_search:
            search_col = global_name_col
            filtered_df = filtered_df[
                filtered_df[search_col].astype(str).str.contains(sku_search, case=False, na=False)]

        st.markdown("### 🎚️ Smart-Scaling & Outlier Control")

        col_out1, col_out2 = st.columns(2)
        with col_out1:
            min_outlier_limit_x = st.number_input(
                f"MINIMUM {x_axis_selector} Limit for Average Calculation (Scope-out Long-tail):",
                value=0.0,
                step=1000000.0 if x_axis_selector == "Gross Sales (Current)" else 1.0,
                help="SKUs below this limit remain VISIBLE, but are EXCLUDED from calculating the X-Axis Average (Center) line."
            )

        with col_out2:
            scale_factor_x = st.number_input(
                f"Right Box Visual Scale (X-Axis Multiplier):",
                value=0.2,
                step=0.1,
                min_value=0.001,
                max_value=10.0,
                help="Set to 1.0 to view the original scale (may look compressed if outliers exist). Set < 1.0 (e.g., 0.2 or 0.1) to shrink the right-side boxes (B7, B8, B9) so the left and middle boxes get more screen space without hiding any SKUs."
            )

        if not filtered_df.empty:
            positive_df = filtered_df[filtered_df[margin_val_col] >= 0]
            total_sales_pos = positive_df['Gross Sales (Current)'].sum()

            if total_sales_pos > 0:
                avg_margin_pct = (positive_df[margin_val_col].sum() / total_sales_pos) * 100
            else:
                avg_margin_pct = positive_df[y_col].mean() if not positive_df.empty else 25.0

            if pd.isna(avg_margin_pct) or avg_margin_pct <= 0:
                avg_margin_pct = 1.0

            avg_y = avg_margin_pct
            def_y_low = avg_y * (2.0 / 3.0)
            def_y_high = avg_y * (4.0 / 3.0)

            normal_x_df = filtered_df[filtered_df[x_col] >= min_outlier_limit_x]
            if not normal_x_df.empty:
                robust_avg_x = normal_x_df[x_col].mean()
            else:
                robust_avg_x = filtered_df[x_col].mean()

            if pd.isna(robust_avg_x) or robust_avg_x <= 0:
                robust_avg_x = 100.0

            def_x_low = robust_avg_x * (2.0 / 3.0)
            def_x_high = robust_avg_x * (4.0 / 3.0)

            x_max_raw = filtered_df[x_col].max()
            x_min_raw = filtered_df[x_col].min()
            x_span_raw = x_max_raw - x_min_raw if x_max_raw != x_min_raw else robust_avg_x

            plot_x_min = x_min_raw - (abs(x_span_raw) * 0.05)
            if x_min_raw < 0:
                plot_x_min = x_min_raw - (abs(x_span_raw) * 0.05)
            plot_x_max = x_max_raw + (abs(x_span_raw) * 0.05)

            y_min_raw = filtered_df[y_col].min()
            y_max_raw = filtered_df[y_col].max()
            y_span_raw = y_max_raw - y_min_raw if y_max_raw != y_min_raw else 100.0

            plot_y_min = y_min_raw - (y_span_raw * 0.05)
            plot_y_max = max(y_max_raw, def_y_high * 1.5) + (y_span_raw * 0.05)
            plot_y_min = min(plot_y_min, def_y_low - (def_y_low * 0.5))
            plot_y_max = max(plot_y_max, def_y_high + (def_y_high * 0.5))
        else:
            def_y_low, def_y_high = 16.67, 33.33
            def_x_low, def_x_high = 66.67, 133.33
            robust_avg_x, avg_margin_pct = 100.0, 25.0
            plot_x_min, plot_x_max, plot_y_min, plot_y_max = -5, 200, -20, 50

        with st.form("threshold_form"):
            st.info(
                "💡 **INFO:** The X and Y axes use the Average value as the absolute center point. The Average lines are marked in Red. Upper and Lower Threshold lines automatically adjust symmetrically against the Average.")

            parent_view_mode = "Common" if view_mode in ["Common (Consolidated Master SKUs)",
                                                         "Commonized Only (Master C- Prefix SKUs)"] else "Uncommon"
            filter_state_str = f"{parent_view_mode}_{market_filter}_{margin_selector}_{x_axis_selector}_{remark_filter}_{status_filter}_{sku_search}_{min_outlier_limit_x}"
            form_key_suffix = hashlib.md5(filter_state_str.encode('utf-8')).hexdigest()

            step_x_input = float(abs(def_x_high - def_x_low) / 2) if def_x_high != def_x_low else 1.0

            col_tx1, col_tx2, col_ty1, col_ty2 = st.columns(4)
            with col_tx1:
                x_low_thresh_input = st.number_input(f"X-Axis Low to Med", value=float(def_x_low), step=step_x_input,
                                                     key=f"xl_{form_key_suffix}")
            with col_tx2:
                x_high_thresh_input = st.number_input(f"X-Axis Med to High", value=float(def_x_high), step=step_x_input,
                                                      key=f"xh_{form_key_suffix}")
            with col_ty1:
                y_low_thresh_input = st.number_input("Y-Axis Low to Med (%)", value=float(def_y_low), step=1.0,
                                                     key=f"yl_{form_key_suffix}")
            with col_ty2:
                y_high_thresh_input = st.number_input("Y-Axis Med to High (%)", value=float(def_y_high), step=1.0,
                                                      key=f"yh_{form_key_suffix}")

            run_thresholds = st.form_submit_button("▶ RUN & UPDATE MATRIX", type="primary")

        y_low_thresh = float(y_low_thresh_input)
        y_high_thresh = float(y_high_thresh_input)
        x_low_thresh = float(x_low_thresh_input)
        x_high_thresh = float(x_high_thresh_input)

        if filtered_df.empty:
            st.error("No data found matching the selected filters.")
        else:
            st.markdown("---")
            st.subheader(f"📈 Matrix View: {market_filter} | {margin_selector} vs {x_axis_selector}")

            x_lbl_low = "Low Growth" if x_axis_selector == "Qty Growth (%)" else "Low Sales"
            x_lbl_med = "Med Growth" if x_axis_selector == "Qty Growth (%)" else "Med Sales"
            x_lbl_high = "High Growth" if x_axis_selector == "Qty Growth (%)" else "High Sales"

            b1_id = f'Box 1 (High Margin, {x_lbl_low})'
            b2_id = f'Box 2 (Med Margin, {x_lbl_low})'
            b3_id = f'Box 3 (Low Margin, {x_lbl_low})'
            b4_id = f'Box 4 (High Margin, {x_lbl_med})'
            b5_id = f'Box 5 (Med Margin, {x_lbl_med})'
            b6_id = f'Box 6 (Low Margin, {x_lbl_med})'
            b7_id = f'Box 7 (High Margin, {x_lbl_high})'
            b8_id = f'Box 8 (Med Margin, {x_lbl_high})'
            b9_id = f'Box 9 (Low Margin, {x_lbl_high})'

            conditions_9box = [
                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] < x_low_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] < x_low_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] < x_low_thresh),

                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] >= x_low_thresh) & (
                        filtered_df[x_col] <= x_high_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] >= x_low_thresh) & (filtered_df[x_col] <= x_high_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] >= x_low_thresh) & (
                        filtered_df[x_col] <= x_high_thresh),

                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] > x_high_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] > x_high_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] > x_high_thresh)
            ]

            choices_9box = [b1_id, b2_id, b3_id, b4_id, b5_id, b6_id, b7_id, b8_id, b9_id]

            filtered_df['Dynamic 9-Box Category'] = np.select(
                conditions_9box, choices_9box, default=b6_id
            )

            if view_mode == "Commonized Only (Master C- Prefix SKUs)":
                filtered_df = filtered_df[filtered_df['New Product Name'].fillna('').astype(str).str.startswith('C-')]
            elif view_mode == "Uncommonized Only (Granular C- Prefix SKUs)":
                if 'New Product Name' in filtered_df.columns:
                    filtered_df = filtered_df[
                        filtered_df['New Product Name'].fillna('').astype(str).str.startswith('C-')]

            render_9box_summary_grid(filtered_df, margin_selector, margin_val_col, y_low_thresh, y_high_thresh,
                                     x_axis_selector, x_low_thresh, x_high_thresh)

            # ==========================================
            # SYNCED BUBBLE CHART RENDERING
            # ==========================================
            plot_df = filtered_df.copy()

            if bubble_size_selector in plot_df.columns:
                plot_df['Bubble_Size'] = pd.to_numeric(plot_df[bubble_size_selector], errors='coerce').fillna(
                    0).abs().replace(0, 1)
            else:
                plot_df['Bubble_Size'] = 10

            box_counts = plot_df['Dynamic 9-Box Category'].value_counts().to_dict()
            total_items = len(plot_df)

            st.info(f"💡 **Total Evaluated & Plotted**: {total_items} SKUs dynamically distributed across the Matrix.")

            name_col = global_name_col

            color_discrete_map = {
                b1_id: '#3871b6', b4_id: '#3871b6', b7_id: '#319b5e',
                b8_id: '#319b5e', b2_id: '#d89f0e', b5_id: '#d89f0e',
                b6_id: '#d89f0e', b9_id: '#d89f0e', b3_id: '#c03d32',
                'Commonized SKU (Black)': '#000000'
            }

            plot_df['Plot_Color_Category'] = plot_df['Dynamic 9-Box Category']
            if 'New Product Name' in plot_df.columns:
                mask_commonized = plot_df['New Product Name'].fillna('').astype(str).str.startswith('C-')
                plot_df.loc[mask_commonized, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

            def format_x_val_str(val):
                if x_axis_selector == "Qty Growth (%)": return f"{val:.1f}%"
                if val >= 1e9 or val <= -1e9: return f"Rp{val / 1e9:,.1f}M"
                return f"Rp{val / 1e6:,.0f}Jt"

            x_format = ':.2f' if x_axis_selector == 'Qty Growth (%)' else ':,.0f'

            def apply_custom_x_scale(v):
                if v <= x_high_thresh:
                    return v
                else:
                    return x_high_thresh + (v - x_high_thresh) * scale_factor_x

            plot_df['X_Plot'] = plot_df[x_col].apply(apply_custom_x_scale)
            plot_df['Y_Plot'] = plot_df[y_col]

            hover_data_dict = {
                'Source_Sheet': True if 'Source_Sheet' in plot_df.columns else False,
                'Dynamic 9-Box Category': True,
                'Plot_Color_Category': False,
                y_col: ':.2f',
                x_col: x_format,
                'X_Plot': False,
                'Y_Plot': False,
                'Bubble_Size': False,
                bubble_size_selector: ':,.0f' if 'Qty' in bubble_size_selector else 'Rp {:,.0f}'
            }

            if 'Qty (Current)' in plot_df.columns and bubble_size_selector != 'Qty (Current)':
                hover_data_dict['Qty (Current)'] = ':,.0f'

            fig = px.scatter(
                plot_df,
                x='X_Plot',
                y='Y_Plot',
                size="Bubble_Size",
                color="Plot_Color_Category",
                color_discrete_map=color_discrete_map,
                hover_name=name_col,
                hover_data=hover_data_dict,
                title=f"9-Box Bubble Chart: {margin_selector} vs {x_axis_selector} (Total: {total_items} SKUs)",
                size_max=60,
                render_mode="svg"
            )

            fig.add_vline(x=apply_custom_x_scale(x_low_thresh), line_dash="dash", line_color="black", line_width=2,
                          opacity=0.8, annotation_text=f" Low ({format_x_val_str(x_low_thresh)})",
                          annotation_position="top left")
            fig.add_vline(x=apply_custom_x_scale(x_high_thresh), line_dash="dash", line_color="black", line_width=2,
                          opacity=0.8, annotation_text=f" High ({format_x_val_str(x_high_thresh)})",
                          annotation_position="top right")
            fig.add_hline(y=y_low_thresh, line_dash="dash", line_color="black", line_width=2, opacity=0.8,
                          annotation_text=f"Low ({y_low_thresh:.1f}%)", annotation_position="bottom right")
            fig.add_hline(y=y_high_thresh, line_dash="dash", line_color="black", line_width=2, opacity=0.8,
                          annotation_text=f"High ({y_high_thresh:.1f}%)", annotation_position="top right")

            fig.add_vline(x=apply_custom_x_scale(robust_avg_x), line_dash="dot", line_color="red", line_width=2,
                          opacity=0.6, annotation_text=f" ← AVG X ({format_x_val_str(robust_avg_x)})",
                          annotation_position="top right", annotation_font_color="red")
            fig.add_hline(y=avg_margin_pct, line_dash="dot", line_color="red", line_width=2, opacity=0.6,
                          annotation_text=f"AVG Y ({avg_margin_pct:.1f}%)", annotation_position="top right",
                          annotation_font_color="red")

            mapped_plot_x_min = apply_custom_x_scale(plot_x_min)
            mapped_plot_x_max = apply_custom_x_scale(plot_x_max)
            mapped_x_low = apply_custom_x_scale(x_low_thresh)
            mapped_x_high = apply_custom_x_scale(x_high_thresh)

            mid_x_low = (mapped_plot_x_min + mapped_x_low) / 2
            mid_x_med = (mapped_x_low + mapped_x_high) / 2
            mid_x_high = (mapped_x_high + mapped_plot_x_max) / 2

            mid_y_low = (plot_y_min + y_low_thresh) / 2
            mid_y_med = (y_low_thresh + y_high_thresh) / 2
            mid_y_high = (y_high_thresh + plot_y_max) / 2

            box_coords = {
                b1_id: {'x': mid_x_low, 'y': mid_y_high, 'id': 'B1'},
                b2_id: {'x': mid_x_low, 'y': mid_y_med, 'id': 'B2'},
                b3_id: {'x': mid_x_low, 'y': mid_y_low, 'id': 'B3'},
                b4_id: {'x': mid_x_med, 'y': mid_y_high, 'id': 'B4'},
                b5_id: {'x': mid_x_med, 'y': mid_y_med, 'id': 'B5'},
                b6_id: {'x': mid_x_med, 'y': mid_y_low, 'id': 'B6'},
                b7_id: {'x': mid_x_high, 'y': mid_y_high, 'id': 'B7'},
                b8_id: {'x': mid_x_high, 'y': mid_y_med, 'id': 'B8'},
                b9_id: {'x': mid_x_high, 'y': mid_y_low, 'id': 'B9'}
            }

            for box_name, coords in box_coords.items():
                count = box_counts.get(box_name, 0)
                if count > 0:
                    is_b5 = (coords['id'] == 'B5')
                    fig.add_annotation(
                        x=coords['x'],
                        y=coords['y'],
                        ax=mapped_x_high + (
                                mapped_plot_x_max - mapped_x_high) * 0.05 if is_b5 else None,
                        ay=mid_y_low if is_b5 else None,
                        axref="x" if is_b5 else None,
                        ayref="y" if is_b5 else None,
                        xref="x",
                        yref="y",
                        text=f"<span style='color:#1f2937;'><b>{coords['id']}</b></span><br><b>{count} SKUs</b>",
                        showarrow=is_b5,
                        arrowhead=2 if is_b5 else 0,
                        arrowsize=1 if is_b5 else 1,
                        arrowwidth=1 if is_b5 else 0.1,
                        arrowcolor="#1f2937" if is_b5 else None,
                        standoff=0,
                        xanchor="left" if is_b5 else "center",
                        font=dict(size=13, color="#374151"),
                        bgcolor="rgba(255, 255, 255, 0.9)" if is_b5 else "rgba(255, 255, 255, 0.8)",
                        bordercolor="rgba(15, 23, 42, 0.8)" if is_b5 else "rgba(15, 23, 42, 0.3)",
                        borderwidth=2 if is_b5 else 1,
                        borderpad=4
                    )

            fig.add_annotation(
                x=0.98,
                y=0.98,
                xref="paper",
                yref="paper",
                text=f"<b>TOTAL:<br>{total_items} SKUs</b>",
                showarrow=False,
                font=dict(size=13, color="white"),
                bgcolor="#374151",
                bordercolor="black",
                borderwidth=1,
                borderpad=5
            )

            ticks_x_actual = [plot_x_min, 0, x_low_thresh, robust_avg_x, x_high_thresh, plot_x_max]
            ticks_x_actual = sorted(list(set([v for v in ticks_x_actual if v >= plot_x_min and v <= plot_x_max])))

            tickvals_x = [apply_custom_x_scale(v) for v in ticks_x_actual]
            ticktext_x = [format_x_val_str(v) for v in ticks_x_actual]

            fig.update_xaxes(
                range=[mapped_plot_x_min, mapped_plot_x_max],
                tickvals=tickvals_x,
                ticktext=ticktext_x,
                title_text="LTM Quantity Growth Rate (%)" if x_axis_selector == "Qty Growth (%)" else "Gross Sales Absolute (IDR)"
            )
            fig.update_yaxes(range=[plot_y_min, plot_y_max], title_text=f"{margin_selector} (%)")

            fig.update_layout(
                height=550,
                hovermode="closest",
                showlegend=True,
                legend=dict(
                    title=None,
                    orientation="h",
                    yanchor="top",
                    y=-0.15,
                    xanchor="center",
                    x=0.5
                ),
                margin=dict(t=50, b=50, l=50, r=50)
            )

            st.plotly_chart(fig, use_container_width=True, config={
                'displayModeBar': True,
                'modeBarButtonsToAdd': ['toImage'],
                'displaylogo': False,
                'toImageButtonOptions': {
                    'format': 'png',
                    'filename': f'9Box_Matrix_{margin_selector.replace(" ", "_")}',
                    'height': 800,
                    'width': 1400,
                    'scale': 2
                }
            })

            # ---------------------------------------------------------
            # DATA DETAIL TABLE & EXPORT
            # ---------------------------------------------------------
            st.markdown("---")
            col_hdr1, col_hdr2, col_hdr3, col_hdr4 = st.columns([1.5, 1, 1, 1])
            with col_hdr1:
                st.subheader("📋 Data Details & Export")

            display_df = filtered_df.drop(columns=['Bubble_Size', 'X_Plot', 'Y_Plot', 'Plot_Color_Category'],
                                          errors='ignore')

            buffer_excel_main = io.BytesIO()
            with pd.ExcelWriter(buffer_excel_main, engine='openpyxl') as writer:
                display_df.to_excel(writer, index=False, sheet_name='9Box_Data')

            with col_hdr2:
                st.download_button(
                    label="📥 Download Excel",
                    data=buffer_excel_main.getvalue(),
                    file_name=f"9Box_Matrix_{margin_selector.replace(' ', '_')}_{market_filter}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    type="primary",
                    use_container_width=True
                )

            with col_hdr3:
                prs = create_portfolio_presentation(filtered_df, margin_selector, margin_val_col)
                buffer_ppt = io.BytesIO()
                prs.save(buffer_ppt)

                st.download_button(
                    label="📥 Download Deck (.pptx)",
                    data=buffer_ppt.getvalue(),
                    file_name=f"Consulting_Deck_{margin_selector.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="secondary",
                    use_container_width=True
                )

            with col_hdr4:
                buffer_html = io.StringIO()
                fig.write_html(buffer_html, include_plotlyjs='cdn')
                st.download_button(
                    label="📥 Download Chart (.html)",
                    data=buffer_html.getvalue(),
                    file_name=f"9Box_Interactive_Chart_{margin_selector.replace(' ', '_')}.html",
                    mime="text/html",
                    type="secondary",
                    use_container_width=True
                )

            format_dict = {}
            if 'Gross Sales (Current)' in display_df.columns: format_dict['Gross Sales (Current)'] = '{:,.0f}'
            if 'Return (Current)' in display_df.columns: format_dict['Return (Current)'] = '{:,.0f}'
            if 'Qty (Current)' in display_df.columns: format_dict['Qty (Current)'] = '{:,.0f}'
            if 'Gross Sales (Previous)' in display_df.columns: format_dict['Gross Sales (Previous)'] = '{:,.0f}'
            if 'Qty (Previous)' in display_df.columns: format_dict['Qty (Previous)'] = '{:,.0f}'

            if 'Gross Profit (Current)' in display_df.columns: format_dict['Gross Profit (Current)'] = '{:,.0f}'
            if 'Gross Profit (%)' in display_df.columns: format_dict['Gross Profit (%)'] = '{:.2f}%'

            if 'Gross Margin (Current)' in display_df.columns: format_dict['Gross Margin (Current)'] = '{:,.0f}'
            if 'Gross Margin (%)' in display_df.columns: format_dict['Gross Margin (%)'] = '{:.2f}%'

            if 'Contribution Margin (Current)' in display_df.columns: format_dict[
                'Contribution Margin (Current)'] = '{:,.0f}'
            if 'Contribution Margin (%)' in display_df.columns: format_dict['Contribution Margin (%)'] = '{:.2f}%'
            if 'Qty Growth (%)' in display_df.columns: format_dict['Qty Growth (%)'] = '{:.2f}%'

            if 'Amount FG' in display_df.columns: format_dict['Amount FG'] = '{:,.0f}'
            if 'Amount Material' in display_df.columns: format_dict['Amount Material'] = '{:,.0f}'

            st.dataframe(
                display_df.style.format(format_dict, na_rep="0.00%"),
                use_container_width=True,
                hide_index=True
            )

            # ---------------------------------------------------------
            # PARETO ANALYSIS MODULE
            # ---------------------------------------------------------
            st.markdown("---")
            col_p1, col_p2, col_p3, col_p4, col_p5 = st.columns([1, 1, 1.2, 1.2, 1.5])

            with col_p1:
                st.subheader("📊 Pareto Analysis")

            with col_p2:
                pareto_market_filter = st.selectbox(
                    "Market Focus:",
                    ["ALL (Local + Export)", "LOCAL", "EXPORT"],
                    key="pareto_market",
                    help="Filter Pareto results based on Market."
                )

            with col_p3:
                chart_margin_label = st.selectbox(
                    "Select Chart Metric:",
                    ["Gross Sales", "Gross Profit", "Gross Margin", "Contribution Margin"],
                    help="SKUs are sorted in descending order based on this metric."
                )

            with col_p4:
                pareto_box_options = ["ALL"] + [f"Box {i}" for i in range(1, 10)]
                pareto_box_filter = st.multiselect(
                    "Filter 9-Box Position:",
                    pareto_box_options,
                    default=["ALL"],
                    help="Select specific quadrants to limit the population included in the Pareto analysis."
                )

            with col_p5:
                pareto_threshold = st.slider(
                    "Pareto Threshold (% SKU)",
                    min_value=1.0,
                    max_value=100.0,
                    value=80.0,
                    step=1.0,
                    help="Calculated from the total GLOBAL population. After securing the Top %, the data is sliced according to the Market Focus selection."
                )

            # --- PARETO LOGIC: CONSISTENCY DATA (POST-FILTERING) ---
            df_for_pareto_base = filtered_df.copy()

            chart_margin_map = {
                "Gross Sales": "Gross Sales (Current)",
                "Gross Profit": "Gross Profit (Current)",
                "Gross Margin": "Gross Margin (Current)",
                "Contribution Margin": "Contribution Margin (Current)"
            }
            chart_margin_col = chart_margin_map[chart_margin_label]

            margin_cols_all = ['Gross Sales (Current)', 'Gross Profit (Current)', 'Gross Margin (Current)',
                               'Contribution Margin (Current)']
            agg_cols = [chart_margin_col]
            for c in ['Qty (Current)'] + margin_cols_all:
                if c in df_for_pareto_base.columns and c not in agg_cols:
                    agg_cols.append(c)

            name_col = global_name_col

            if chart_margin_col in df_for_pareto_base.columns:
                df_pareto_global = df_for_pareto_base.groupby(name_col)[agg_cols].sum().reset_index()
                df_pareto_global = df_pareto_global.sort_values(by=chart_margin_col, ascending=False).reset_index(
                    drop=True)

                total_all_skus_global = len(df_pareto_global)
                if pareto_threshold >= 100:
                    cutoff_count_global = total_all_skus_global
                else:
                    cutoff_count_global = max(1, round(total_all_skus_global * pareto_threshold / 100))

                top_global_skus = df_pareto_global.iloc[:cutoff_count_global][name_col].tolist()

                df_sliced = df_for_pareto_base.copy()

                if pareto_market_filter != "ALL (Local + Export)":
                    if 'Source_Sheet' in df_sliced.columns:
                        df_sliced = df_sliced[df_sliced['Source_Sheet'] == pareto_market_filter]

                if "ALL" not in pareto_box_filter and len(pareto_box_filter) > 0:
                    box_pattern = '|'.join([f"^{b}" for b in pareto_box_filter])
                    df_sliced = df_sliced[
                        df_sliced['Dynamic 9-Box Category'].str.contains(box_pattern, regex=True, na=False)]

                total_sliced_skus_base = df_sliced[name_col].nunique()
                base_sliced_sales = df_sliced[
                    'Gross Sales (Current)'].sum() if 'Gross Sales (Current)' in df_sliced.columns else 0
                base_sliced_qty = df_sliced[
                    'Qty (Current)'].sum() if 'Qty (Current)' in df_sliced.columns else 0

                # Only retrieve data from slicer results that belong to Top Global SKUs
                df_pareto_raw_filtered = df_sliced[df_sliced[name_col].isin(top_global_skus)].copy()

                # 3. FINAL GROUPING FOR DISPLAY
                df_pareto_filtered = df_pareto_raw_filtered.groupby(name_col)[agg_cols].sum().reset_index()
                df_pareto_filtered = df_pareto_filtered.sort_values(by=chart_margin_col, ascending=False).reset_index(
                    drop=True)

                if 'Dynamic 9-Box Category' in df_for_pareto_base.columns:
                    box_category_map = df_for_pareto_base.groupby(name_col)['Dynamic 9-Box Category'].agg(
                        lambda s: s.iloc[0] if s.nunique() == 1 else f"Mixed ({s.nunique()} box)"
                    )
                    df_pareto_filtered['9-Box Kwadran'] = df_pareto_filtered[name_col].map(box_category_map)

                if not df_pareto_filtered.empty:
                    expanded_cols_to_add = [c for c in df_for_pareto_base.columns if
                                            c.startswith('Produk_Key_') or c.startswith('Product Name_')]
                    if expanded_cols_to_add:
                        mapping_df = df_for_pareto_base[[name_col] + expanded_cols_to_add].drop_duplicates(
                            subset=[name_col])
                        df_pareto_filtered = pd.merge(df_pareto_filtered, mapping_df, on=name_col, how='left')

                    total_val = df_pareto_filtered[chart_margin_col].sum()
                    if total_val != 0:
                        df_pareto_filtered['Cumulative %'] = (df_pareto_filtered[
                                                                  chart_margin_col].cumsum() / total_val) * 100
                    else:
                        df_pareto_filtered['Cumulative %'] = 0.0

                    total_pareto_skus = len(df_pareto_filtered)

                    summary_sku = total_pareto_skus
                    summary_qty = df_pareto_filtered[
                        'Qty (Current)'].sum() if 'Qty (Current)' in df_pareto_filtered.columns else 0
                    summary_sales = df_pareto_filtered[
                        'Gross Sales (Current)'].sum() if 'Gross Sales (Current)' in df_pareto_filtered.columns else 0

                    pct_sku = (summary_sku / total_sliced_skus_base * 100) if total_sliced_skus_base > 0 else 0
                    pct_sales = (summary_sales / base_sliced_sales * 100) if base_sliced_sales > 0 else 0
                    pct_qty = (summary_qty / base_sliced_qty * 100) if base_sliced_qty > 0 else 0

                    m1, m2, m3 = st.columns(3)

                    if pareto_market_filter == "ALL (Local + Export)" and "ALL" in pareto_box_filter:
                        txt_sku = f"↑ {pareto_threshold:.1f}% from {total_all_skus_global:,} Global SKUs"
                    else:
                        txt_sku = f"↑ {summary_sku:,} Sliced SKUs (from {cutoff_count_global:,} Top Global)"

                    m1.metric("Total SKUs", f"{summary_sku:,}", txt_sku)
                    m2.metric("Quantity", f"{summary_qty:,.0f} ({pct_qty:.1f}%)", f"↑ {pct_qty:.1f}% from total slice")

                    summary_sales_bn = summary_sales / 1e9
                    m3.metric("Gross Sales", f"Rp {summary_sales_bn:,.1f} Bn ({pct_sales:.1f}%)",
                              f"↑ {pct_sales:.1f}% from total slice")

                    margin_label_map = {
                        'Gross Profit (Current)': 'Gross Profit',
                        'Gross Margin (Current)': 'Gross Margin',
                        'Contribution Margin (Current)': 'Contribution Margin'
                    }

                    mg1, mg2, mg3 = st.columns(3)

                    loop_margins = ['Gross Profit (Current)', 'Gross Margin (Current)', 'Contribution Margin (Current)']
                    for col_widget, margin_col in zip([mg1, mg2, mg3], loop_margins):
                        if margin_col not in df_pareto_filtered.columns:
                            continue
                        margin_sum = df_pareto_filtered[margin_col].sum()
                        grand_total = df_sliced[margin_col].sum() if margin_col in df_sliced.columns else 0
                        cumulative_pct = (margin_sum / grand_total * 100) if grand_total > 0 else 0
                        margin_ratio_pct = (margin_sum / summary_sales * 100) if summary_sales > 0 else 0
                        label = margin_label_map[margin_col]
                        if margin_col == chart_margin_col:
                            label = f"📊 {label} (in Chart)"
                        margin_sum_bn = margin_sum / 1e9
                        col_widget.metric(
                            label,
                            f"Rp {margin_sum_bn:,.1f} Bn ({cumulative_pct:.1f}%)",
                            f"{margin_ratio_pct:.1f}% margin ratio"
                        )

                    fig_pareto = go.Figure()

                    bar_colors = []
                    for i, row in df_pareto_filtered.iterrows():
                        npm = ""
                        if 'New Product Name' in df_for_pareto_base.columns:
                            matches = df_for_pareto_base[df_for_pareto_base[name_col] == row[name_col]]
                            if not matches.empty:
                                npm = str(matches.iloc[0]['New Product Name'])

                        if npm.startswith('C-'):
                            bar_colors.append('#000000')
                        elif row[chart_margin_col] < 0:
                            bar_colors.append('#c03d32')
                        else:
                            bar_colors.append('#38bdf8')

                    fig_pareto.add_trace(go.Bar(
                        x=df_pareto_filtered[name_col],
                        y=df_pareto_filtered[chart_margin_col],
                        name=chart_margin_label,
                        marker_color=bar_colors,
                        yaxis='y1'
                    ))

                    fig_pareto.add_trace(go.Scatter(
                        x=df_pareto_filtered[name_col],
                        y=df_pareto_filtered['Cumulative %'],
                        name=f'Cumulative % ({chart_margin_label})',
                        marker_color='#f97316',
                        mode='lines+markers',
                        yaxis='y2'
                    ))

                    fig_pareto.update_layout(
                        title=f"Pareto Chart: {chart_margin_label} (Sliced from Top {pareto_threshold}% Global | {total_pareto_skus} SKUs)",
                        hovermode="x unified",
                        height=550,
                        xaxis=dict(showticklabels=False, title=f"SKUs (Ranked by {chart_margin_label})"),
                        yaxis=dict(title=f"{chart_margin_label} (IDR)"),
                        yaxis2=dict(
                            title=f"Cumulative % {chart_margin_label}",
                            overlaying='y',
                            side='right',
                            showgrid=False
                        ),
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
                    )

                    st.plotly_chart(fig_pareto, use_container_width=True)

                    # --- DUAL EXPORT BUTTONS FOR COMMON MODE ---
                    col_pt_dl1, col_pt_dl2 = st.columns(2)

                    # 1. Download Master SKU Pareto (As is)
                    buffer_excel_pareto = io.BytesIO()
                    with pd.ExcelWriter(buffer_excel_pareto, engine='openpyxl') as writer:
                        df_pareto_filtered.to_excel(writer, index=False, sheet_name='Pareto_Data')

                    with col_pt_dl1:
                        st.download_button(
                            label=f"📥 Download Top {pareto_threshold}% Master SKUs (Excel)",
                            data=buffer_excel_pareto.getvalue(),
                            file_name=f"Pareto_Analysis_{chart_margin_label.replace(' ', '_')}_{pareto_threshold}pctSKU.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            type="primary",
                            use_container_width=True
                        )

                    # 2. Download Exploded Granular Pareto (Pulling Raw Data for Pure Numbers)
                    if view_mode in ["Common (Consolidated Master SKUs)", "Commonized Only (Master C- Prefix SKUs)"]:
                        pareto_master_skus = df_pareto_filtered[name_col].tolist()

                        df_pareto_exploded = df_raw[df_raw['New Product Name'].isin(pareto_master_skus)].copy()

                        if pareto_market_filter != "ALL (Local + Export)":
                            if 'Source_Sheet' in df_pareto_exploded.columns:
                                df_pareto_exploded = df_pareto_exploded[
                                    df_pareto_exploded['Source_Sheet'] == pareto_market_filter]

                        if '9-Box Kwadran' in df_pareto_filtered.columns:
                            mapping_box = df_pareto_filtered.set_index(name_col)['9-Box Kwadran'].to_dict()
                            df_pareto_exploded['Master 9-Box Kwadran'] = df_pareto_exploded['New Product Name'].map(
                                mapping_box)

                        front_cols = ['Produk_Key', 'Product Name', 'New Code', 'New Product Name',
                                      'Master 9-Box Kwadran']
                        front_cols = [c for c in front_cols if c in df_pareto_exploded.columns]
                        final_cols = front_cols + [c for c in df_pareto_exploded.columns if c not in front_cols]
                        df_pareto_exploded = df_pareto_exploded[final_cols]

                        if 'Gross Sales (Current)' in df_pareto_exploded.columns:
                            df_pareto_exploded = df_pareto_exploded.sort_values(
                                by=['New Product Name', 'Gross Sales (Current)'], ascending=[True, False])

                        buffer_excel_exploded = io.BytesIO()
                        with pd.ExcelWriter(buffer_excel_exploded, engine='openpyxl') as writer:
                            df_pareto_exploded.to_excel(writer, index=False, sheet_name='Granular_Pareto')

                        with col_pt_dl2:
                            st.download_button(
                                label=f"📥 Download Top {pareto_threshold}% by Old Product_Key (Excel)",
                                data=buffer_excel_exploded.getvalue(),
                                file_name=f"Pareto_Granular_OldKeys_{chart_margin_label.replace(' ', '_')}_{pareto_threshold}pct.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                type="secondary",
                                use_container_width=True
                            )

                    st.dataframe(
                        df_pareto_filtered.style.format({
                            'Gross Sales (Current)': '{:,.0f}',
                            'Qty (Current)': '{:,.0f}',
                            'Gross Profit (Current)': '{:,.0f}',
                            'Gross Margin (Current)': '{:,.0f}',
                            'Contribution Margin (Current)': '{:,.0f}',
                            'Cumulative %': '{:.2f}%'
                        }),
                        use_container_width=True,
                        hide_index=True
                    )

                    st.markdown("---")
                    st.subheader(f"🧩 9-Box Summary for Top {pareto_threshold}% Pareto SKUs")
                    st.info(
                        "The visualizations below (Matrix Grid & Bubble Chart) display the 9-box distribution exclusively for the SKUs included in the Pareto filter above.")

                    pareto_sku_list = df_pareto_filtered[name_col].tolist()
                    df_pareto_9box = df_pareto_raw_filtered.copy()

                    if not df_pareto_9box.empty:
                        tab_p_matrix, tab_p_bubble = st.tabs(["🧮 Matrix Grid", "🫧 Bubble Chart"])

                        with tab_p_matrix:
                            html_matrix_pareto = render_9box_summary_grid(
                                df_pareto_9box,
                                margin_selector,
                                margin_val_col,
                                y_low_thresh,
                                y_high_thresh,
                                x_axis_selector,
                                x_low_thresh,
                                x_high_thresh
                            )

                            st.markdown("##### 📥 Export Pareto 9-Box Data")
                            col_dl_mat1, col_dl_mat2 = st.columns(2)
                            with col_dl_mat1:
                                st.download_button(
                                    label="📥 Download Matrix Grid (.html)",
                                    data=html_matrix_pareto,
                                    file_name=f"Pareto_9Box_Matrix_{margin_selector.replace(' ', '_')}.html",
                                    mime="text/html",
                                    type="secondary",
                                    use_container_width=True
                                )
                            with col_dl_mat2:
                                excel_pareto_9box = get_raw_excel(df_pareto_9box)
                                st.download_button(
                                    label=f"📥 Download 9-Box Data for Top {pareto_threshold}% (Excel)",
                                    data=excel_pareto_9box,
                                    file_name=f"Pareto_9Box_Top{pareto_threshold}_{chart_margin_label.replace(' ', '')}.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                    type="primary",
                                    use_container_width=True
                                )

                        with tab_p_bubble:
                            plot_df_p = df_pareto_9box.copy()

                            if bubble_size_selector in plot_df_p.columns:
                                plot_df_p['Bubble_Size'] = pd.to_numeric(plot_df_p[bubble_size_selector],
                                                                         errors='coerce').fillna(0).abs().replace(0, 1)
                            else:
                                plot_df_p['Bubble_Size'] = 10

                            box_counts_p = plot_df_p['Dynamic 9-Box Category'].value_counts().to_dict()
                            total_items_p = len(plot_df_p)

                            plot_df_p['X_Plot'] = plot_df_p[x_col].apply(apply_custom_x_scale)
                            plot_df_p['Y_Plot'] = plot_df_p[y_col]

                            plot_df_p['Plot_Color_Category'] = plot_df_p['Dynamic 9-Box Category']
                            if 'New Product Name' in plot_df_p.columns:
                                mask_commonized_p = plot_df_p['New Product Name'].fillna('').astype(str).str.startswith(
                                    'C-')
                                plot_df_p.loc[mask_commonized_p, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

                            hover_data_dict_p = hover_data_dict.copy()

                            fig_p = px.scatter(
                                plot_df_p,
                                x='X_Plot',
                                y='Y_Plot',
                                size="Bubble_Size",
                                color="Plot_Color_Category",
                                color_discrete_map=color_discrete_map,
                                hover_name=name_col,
                                hover_data=hover_data_dict_p,
                                title=f"Pareto 9-Box Bubble Chart: {margin_selector} vs {x_axis_selector} (Total: {total_items_p} SKUs)",
                                size_max=60,
                                render_mode="svg"
                            )

                            fig_p.add_vline(x=apply_custom_x_scale(x_low_thresh), line_dash="dash", line_color="black",
                                            line_width=2, opacity=0.8,
                                            annotation_text=f" Low ({format_x_val_str(x_low_thresh)})",
                                            annotation_position="top left")
                            fig_p.add_vline(x=apply_custom_x_scale(x_high_thresh), line_dash="dash", line_color="black",
                                            line_width=2, opacity=0.8,
                                            annotation_text=f" High ({format_x_val_str(x_high_thresh)})",
                                            annotation_position="top right")
                            fig_p.add_hline(y=y_low_thresh, line_dash="dash", line_color="black", line_width=2,
                                            opacity=0.8, annotation_text=f"Low ({y_low_thresh:.1f}%)",
                                            annotation_position="bottom right")
                            fig_p.add_hline(y=y_high_thresh, line_dash="dash", line_color="black", line_width=2,
                                            opacity=0.8, annotation_text=f"High ({y_high_thresh:.1f}%)",
                                            annotation_position="top right")

                            fig_p.add_vline(x=apply_custom_x_scale(robust_avg_x), line_dash="dot", line_color="red",
                                            line_width=2, opacity=0.6,
                                            annotation_text=f" ← AVG X ({format_x_val_str(robust_avg_x)})",
                                            annotation_position="top right", annotation_font_color="red")
                            fig_p.add_hline(y=avg_margin_pct, line_dash="dot", line_color="red", line_width=2,
                                            opacity=0.6, annotation_text=f"AVG Y ({avg_margin_pct:.1f}%)",
                                            annotation_position="top right", annotation_font_color="red")

                            for box_name, coords in box_coords.items():
                                count = box_counts_p.get(box_name, 0)
                                if count > 0:
                                    is_b5 = (coords['id'] == 'B5')
                                    fig_p.add_annotation(
                                        x=coords['x'],
                                        y=coords['y'],
                                        ax=mapped_x_high + (
                                                mapped_plot_x_max - mapped_x_high) * 0.05 if is_b5 else None,
                                        ay=mid_y_low if is_b5 else None,
                                        axref="x" if is_b5 else None,
                                        ayref="y" if is_b5 else None,
                                        xref="x",
                                        yref="y",
                                        text=f"<span style='color:#1f2937;'><b>{coords['id']}</b></span><br><b>{count} SKUs</b>",
                                        showarrow=is_b5,
                                        arrowhead=2 if is_b5 else 0,
                                        arrowsize=1 if is_b5 else 1,
                                        arrowwidth=1 if is_b5 else 0.1,
                                        arrowcolor="#1f2937" if is_b5 else None,
                                        standoff=0,
                                        xanchor="left" if is_b5 else "center",
                                        font=dict(size=13, color="#374151"),
                                        bgcolor="rgba(255, 255, 255, 0.9)" if is_b5 else "rgba(255, 255, 255, 0.8)",
                                        bordercolor="rgba(15, 23, 42, 0.8)" if is_b5 else "rgba(15, 23, 42, 0.3)",
                                        borderwidth=2 if is_b5 else 1,
                                        borderpad=4
                                    )

                            fig_p.add_annotation(
                                x=0.98,
                                y=0.98,
                                xref="paper",
                                yref="paper",
                                text=f"<b>TOTAL:<br>{total_items_p} SKUs</b>",
                                showarrow=False,
                                font=dict(size=13, color="white"),
                                bgcolor="#374151",
                                bordercolor="black",
                                borderwidth=1,
                                borderpad=5
                            )

                            fig_p.update_xaxes(
                                range=[mapped_plot_x_min, mapped_plot_x_max],
                                tickvals=tickvals_x,
                                ticktext=ticktext_x,
                                title_text="LTM Quantity Growth Rate (%)" if x_axis_selector == "Qty Growth (%)" else "Gross Sales Absolute (IDR)"
                            )
                            fig_p.update_yaxes(range=[plot_y_min, plot_y_max], title_text=f"{margin_selector} (%)")

                            fig_p.update_layout(
                                height=550,
                                hovermode="closest",
                                showlegend=True,
                                legend=dict(
                                    title=None,
                                    orientation="h",
                                    yanchor="top",
                                    y=-0.15,
                                    xanchor="center",
                                    x=0.5
                                ),
                                margin=dict(t=50, b=50, l=50, r=50)
                            )

                            st.plotly_chart(fig_p, use_container_width=True, config={
                                'displayModeBar': True,
                                'modeBarButtonsToAdd': ['toImage'],
                                'displaylogo': False,
                                'toImageButtonOptions': {
                                    'format': 'png',
                                    'filename': f'Pareto_9Box_Bubble_{margin_selector.replace(" ", "_")}',
                                    'height': 800,
                                    'width': 1400,
                                    'scale': 2
                                }
                            })

                            buffer_html_p = io.StringIO()
                            fig_p.write_html(buffer_html_p, include_plotlyjs='cdn')
                            st.download_button(
                                label="📥 Download Bubble Chart (.html)",
                                data=buffer_html_p.getvalue(),
                                file_name=f"Pareto_9Box_Bubble_{margin_selector.replace(' ', '_')}.html",
                                mime="text/html",
                                type="secondary"
                            )
                    else:
                        st.warning("No valid data available to render the 9-Box Diagram.")

                else:
                    st.warning(f"No data found for {chart_margin_label} to generate a Pareto Chart.")
            else:
                st.error(f"Metric '{chart_margin_label}' not available in the dataset.")

    # ---------------------------------------------------------
    # TAB 2: VENN DIAGRAM INFOGRAPHIC (4-LEAF CLOVER)
    # ---------------------------------------------------------
    with tab_b3_intersection:
        st.markdown("### 📊 Strategic Portfolio Intersection (4-Set Venn Analysis)")
        st.info(
            "💡 **Executive Infographic (4-Leaf Clover):** Discover precise intersections of the 4 Main Criteria. Select **'None'** to exclude a circle from the intersection analysis.")

        with st.form("venn_settings_form"):
            st.markdown("#### 🎯 Global & Venn Filters")

            col_g1, col_g2, col_g3 = st.columns(3)
            with col_g1:
                v_market = st.selectbox("Global Filter: Market Focus", ["ALL", "LOCAL", "EXPORT"],
                                        help="Filter the initial population based on Market.", key="venn_market")
            with col_g2:
                v_margin_selector = st.selectbox("Y-Axis Metric (Margin):",
                                                 ["Gross Profit", "Gross Margin", "Contribution Margin"],
                                                 key="venn_margin")
            with col_g3:
                v_x_axis_selector = st.selectbox("X-Axis Metric (Performance):",
                                                 ["Gross Sales (Current)", "Qty Growth (%)"], key="venn_xaxis")

            col_v1, col_v2, col_v3, col_v4 = st.columns(4)

            with col_v1:
                v_newcode = st.selectbox("Set A: Commonization", ["ALL", "None", "Commonized SKU", "Others"],
                                         help="SKU Type based on 'C-' prefix.")

            with col_v2:
                if 'Remark' in df_main.columns:
                    valid_rem1 = sorted(
                        [str(x) for x in df_main['Remark'].unique() if pd.notna(x) and str(x).strip() != ''])
                    rem1_options = ["ALL", "None"] + valid_rem1
                else:
                    rem1_options = ["ALL", "None"]
                v_remark = st.multiselect("Set B: Plan Discontinue", rem1_options, default=["ALL"])

            with col_v3:
                if 'Remark2' in df_main.columns:
                    valid_rem2 = sorted(
                        [str(x) for x in df_main['Remark2'].unique() if pd.notna(x) and str(x).strip() != ''])
                    rem2_options = ["ALL", "None"] + valid_rem2 + ["Kosong (Blank)"]
                else:
                    rem2_options = ["ALL", "None", "McK", "Kosong (Blank)"]
                v_remark2 = st.selectbox("Set C: McKinsey Project", rem2_options)

            with col_v4:
                box_list = ["ALL", "None"] + [f"Box {i}" for i in range(1, 10)]
                v_box = st.multiselect("Set D: 9-Box Position", box_list, default=["ALL"])

            run_venn = st.form_submit_button("▶ GENERATE 4-LEAF CLOVER INFOGRAPHIC", type="primary")

        # --- VENN FILTER LOGIC ---
        df_v = df_main.copy()
        if v_market != "ALL":
            df_v = df_v[df_v['Source_Sheet'] == v_market]

        if v_margin_selector == "Gross Profit":
            y_col_v = 'Gross Profit (%)'
            margin_val_col_v = 'Gross Profit (Current)'
        elif v_margin_selector == "Gross Margin":
            y_col_v = 'Gross Margin (%)'
            margin_val_col_v = 'Gross Margin (Current)'
        else:
            y_col_v = 'Contribution Margin (%)'
            margin_val_col_v = 'Contribution Margin (Current)'

        x_col_v = v_x_axis_selector

        if not df_v.empty:
            positive_df_v = df_v[df_v[margin_val_col_v] >= 0]
            total_sales_pos_v = positive_df_v['Gross Sales (Current)'].sum() if not positive_df_v.empty else 0

            if total_sales_pos_v > 0:
                avg_margin_pct_v = (positive_df_v[margin_val_col_v].sum() / total_sales_pos_v) * 100
            else:
                avg_margin_pct_v = positive_df_v[y_col_v].mean() if not positive_df_v.empty else 25.0

            if pd.isna(avg_margin_pct_v) or avg_margin_pct_v <= 0: avg_margin_pct_v = 1.0

            v_y_low = avg_margin_pct_v * (2.0 / 3.0)
            v_y_high = avg_margin_pct_v * (4.0 / 3.0)

            avg_x_v = df_v[x_col_v].mean() if not df_v.empty else 100.0
            if pd.isna(avg_x_v) or avg_x_v <= 0: avg_x_v = 100.0
            v_x_low = avg_x_v * (2.0 / 3.0)
            v_x_high = avg_x_v * (4.0 / 3.0)

            c_box_v = [
                (df_v[y_col_v] > v_y_high) & (df_v[x_col_v] < v_x_low),
                (df_v[y_col_v] >= v_y_low) & (df_v[y_col_v] <= v_y_high) & (df_v[x_col_v] < v_x_low),
                (df_v[y_col_v] < v_y_low) & (df_v[x_col_v] < v_x_low),
                (df_v[y_col_v] > v_y_high) & (df_v[x_col_v] >= v_x_low) & (df_v[x_col_v] <= v_x_high),
                (df_v[y_col_v] >= v_y_low) & (df_v[y_col_v] <= v_y_high) & (df_v[x_col_v] >= v_x_low) & (
                        df_v[x_col_v] <= v_x_high),
                (df_v[y_col_v] < v_y_low) & (df_v[x_col_v] >= v_x_low) & (df_v[x_col_v] <= v_x_high),
                (df_v[y_col_v] > v_y_high) & (df_v[x_col_v] > v_x_high),
                (df_v[y_col_v] >= v_y_low) & (df_v[y_col_v] <= v_y_high) & (df_v[x_col_v] > v_x_high),
                (df_v[y_col_v] < v_y_low) & (df_v[x_col_v] > v_x_high)
            ]
            choices_v = [f"Box {i}" for i in [1, 2, 3, 4, 5, 6, 7, 8, 9]]
            df_v['Dynamic 9-Box Category'] = np.select(c_box_v, choices_v, default="Box 6")
            df_v['Venn_Box'] = df_v['Dynamic 9-Box Category']

            # SET A (Commonization)
            setA = pd.Series(True, index=df_v.index)
            if v_newcode == "None":
                setA = pd.Series(False, index=df_v.index)
            elif v_newcode == "Commonized SKU":
                setA = df_v['New Product Name'].fillna('').astype(str).str.startswith('C-')
            elif v_newcode == "Others":
                setA = ~df_v['New Product Name'].fillna('').astype(str).str.startswith('C-')

            # SET B (Plan Discontinue / Remark)
            setB = pd.Series(True, index=df_v.index)
            if "None" in v_remark:
                setB = pd.Series(False, index=df_v.index)
            else:
                actual_remarks_v = [r for r in v_remark if r not in ["ALL", "None"]]
                if actual_remarks_v:
                    setB = df_v['Remark'].astype(str).isin(actual_remarks_v)

            # SET C (McKinsey Project / Remark2)
            setC = pd.Series(True, index=df_v.index)
            if v_remark2 == "None":
                setC = pd.Series(False, index=df_v.index)
            elif v_remark2 == "Kosong (Blank)":
                setC = (df_v['Remark2'].astype(str).str.strip() == '')
            elif v_remark2 != "ALL":
                setC = (df_v['Remark2'].astype(str).str.strip().str.upper() == v_remark2.upper())

            # SET D (9-Box Position)
            setD = pd.Series(True, index=df_v.index)
            if "None" in v_box:
                setD = pd.Series(False, index=df_v.index)
            else:
                if "ALL" not in v_box and len(v_box) > 0:
                    setD = df_v['Venn_Box'].isin(v_box)

            # --- INTERSECTION MASKS (V1 - V15) ---
            m_A_only = setA & ~setB & ~setC & ~setD
            m_B_only = ~setA & setB & ~setC & ~setD
            m_C_only = ~setA & ~setB & setC & ~setD
            m_D_only = ~setA & ~setB & ~setC & setD

            m_AB_only = setA & setB & ~setC & ~setD
            m_AC_only = setA & ~setB & setC & ~setD
            m_AD_only = setA & ~setB & ~setC & setD
            m_BC_only = ~setA & setB & setC & ~setD
            m_BD_only = ~setA & setB & ~setC & setD
            m_CD_only = ~setA & ~setB & setC & setD

            m_ABC_only = setA & setB & setC & ~setD
            m_ABD_only = setA & setB & ~setC & setD
            m_ACD_only = setA & ~setB & setC & setD
            m_BCD_only = ~setA & setB & setC & setD

            m_ABCD = setA & setB & setC & setD

            def get_venn_label(bidang, mask):
                cnt = mask.sum()
                if cnt == 0: return ""
                sales = df_v.loc[mask, 'Gross Sales (Current)'].sum()
                return f"<b>{bidang}</b><br>{cnt} SKUs<br>Rp {sales / 1e9:,.1f} Bn"

            st.markdown("---")
            st.markdown("#### 🎯 4-Leaf Clover Venn Diagram Infographic")

            # VENN RENDERING (ELLIPSE/CIRCLES)
            fig_venn = go.Figure()

            # Dynamic Transparent Colors Logic
            c_A_fill = "rgba(0,0,0,0)" if v_newcode == "None" else "rgba(15, 118, 110, 0.3)"
            c_A_line = "rgba(0,0,0,0)" if v_newcode == "None" else "#0f766e"

            c_B_fill = "rgba(0,0,0,0)" if "None" in v_remark else "rgba(217, 119, 6, 0.3)"
            c_B_line = "rgba(0,0,0,0)" if "None" in v_remark else "#d97706"

            c_C_fill = "rgba(0,0,0,0)" if v_remark2 == "None" else "rgba(30, 58, 138, 0.3)"
            c_C_line = "rgba(0,0,0,0)" if v_remark2 == "None" else "#1e3a8a"

            c_D_fill = "rgba(0,0,0,0)" if "None" in v_box else "rgba(22, 163, 74, 0.3)"
            c_D_line = "rgba(0,0,0,0)" if "None" in v_box else "#16a34a"

            # Circles Rendering
            fig_venn.add_shape(type="circle", x0=1.5, y0=3.5, x1=6.5, y1=8.5, fillcolor=c_A_fill, line_color=c_A_line,
                               line_width=3)
            fig_venn.add_shape(type="circle", x0=3.5, y0=3.5, x1=8.5, y1=8.5, fillcolor=c_B_fill, line_color=c_B_line,
                               line_width=3)
            fig_venn.add_shape(type="circle", x0=3.5, y0=1.5, x1=8.5, y1=6.5, fillcolor=c_C_fill, line_color=c_C_line,
                               line_width=3)
            fig_venn.add_shape(type="circle", x0=1.5, y0=1.5, x1=6.5, y1=6.5, fillcolor=c_D_fill, line_color=c_D_line,
                               line_width=3)

            # Title Formatting
            lbl_A_text = "None" if v_newcode == "None" else v_newcode
            lbl_B_text = "None" if "None" in v_remark else (', '.join(v_remark) if 'ALL' not in v_remark else 'ALL')
            lbl_C_text = "None" if v_remark2 == "None" else v_remark2
            lbl_D_text = "None" if "None" in v_box else (', '.join(v_box) if 'ALL' not in v_box else 'ALL')

            lbl_A = f"Commonization<br>({lbl_A_text})"
            lbl_B = f"Plan Discontinue<br>({lbl_B_text})"
            lbl_C = f"McKinsey Project<br>({lbl_C_text})"
            lbl_D = f"9-Box Position<br>({lbl_D_text})"

            # Title Injection
            if v_newcode != "None":
                fig_venn.add_annotation(x=2.5, y=9.2, text=f"<b>{lbl_A}</b>", showarrow=False,
                                        font=dict(size=14, color="#0f766e"))
            if "None" not in v_remark:
                fig_venn.add_annotation(x=7.5, y=9.2, text=f"<b>{lbl_B}</b>", showarrow=False,
                                        font=dict(size=14, color="#d97706"))
            if v_remark2 != "None":
                fig_venn.add_annotation(x=7.5, y=0.8, text=f"<b>{lbl_C}</b>", showarrow=False,
                                        font=dict(size=14, color="#1e3a8a"))
            if "None" not in v_box:
                fig_venn.add_annotation(x=2.5, y=0.8, text=f"<b>{lbl_D}</b>", showarrow=False,
                                        font=dict(size=14, color="#16a34a"))

            # V1 to V15 Labels (Clover Layout)
            annotations = [
                (2.5, 7.5, "V1", m_A_only),
                (7.5, 7.5, "V2", m_B_only),
                (7.5, 2.5, "V3", m_C_only),
                (2.5, 2.5, "V4", m_D_only),
                (5.0, 7.8, "V5", m_AB_only),
                (2.2, 5.0, "V7", m_AD_only),
                (7.8, 5.0, "V8", m_BC_only),
                (5.0, 2.2, "V10", m_CD_only),
                (5.8, 5.8, "V11", m_ABC_only),
                (4.2, 5.8, "V12", m_ABD_only),
                (4.2, 4.2, "V13", m_ACD_only),
                (5.8, 4.2, "V14", m_BCD_only),
            ]

            for ax_x, ax_y, bid, m in annotations:
                txt = get_venn_label(bid, m)
                if txt:
                    fig_venn.add_annotation(x=ax_x, y=ax_y, text=txt, showarrow=False,
                                            font=dict(size=12, color="black"))

            bullseye_label = get_venn_label('', m_ABCD)
            if bullseye_label != "":
                fig_venn.add_annotation(x=5.0, y=5.0,
                                        text=f"<span style='color:red;'><b>V15<br>(BULLSEYE)</b></span><br>{bullseye_label.replace('<b></b><br>', '')}",
                                        showarrow=False, font=dict(size=13, color="black"))

            # FIX FOR OVERLAPPING V6 & V9 AREAS: Using Pointer Lines
            v6_text = get_venn_label("V6", m_AC_only)
            if v6_text != "":
                fig_venn.add_annotation(
                    x=5.0, y=5.0, ax=1.5, ay=5.0, axref="x", ayref="y",
                    text=v6_text, showarrow=True, arrowhead=2, arrowsize=1, arrowwidth=1.5, arrowcolor="#0f766e",
                    font=dict(size=12, color="black")
                )

            v9_text = get_venn_label("V9", m_BD_only)
            if v9_text != "":
                fig_venn.add_annotation(
                    x=5.0, y=5.0, ax=8.5, ay=5.0, axref="x", ayref="y",
                    text=v9_text, showarrow=True, arrowhead=2, arrowsize=1, arrowwidth=1.5, arrowcolor="#d97706",
                    font=dict(size=12, color="black")
                )

            fig_venn.update_xaxes(visible=False, range=[0, 10])
            fig_venn.update_yaxes(visible=False, range=[0, 10])
            fig_venn.update_layout(height=650, plot_bgcolor='white', margin=dict(t=30, b=30, l=30, r=30))

            st.plotly_chart(fig_venn, use_container_width=True)
            st.caption(
                "*Visual Note: The diagonal intersections V6 (A∩C) and V9 (B∩D), which are geometrically hidden, are now precisely indicated using pointer lines to the center, ensuring the total SKU count on the chart represents exactly 100% of your manual calculations.*")

            # --- EXPORT TABLES V1 - V15 ---
            st.markdown("### 📥 Download Area (V1 - V15)")
            st.info("Download pure SKU data per Venn Diagram sector for further calculation in Excel.")

            export_grids = st.columns(4)
            bidang_data = [
                ("V1 (Set A Only)", m_A_only),
                ("V2 (Set B Only)", m_B_only),
                ("V3 (Set C Only)", m_C_only),
                ("V4 (Set D Only)", m_D_only),
                ("V5 (V1 ∩ V2)", m_AB_only),
                ("V6 (V1 ∩ V3)", m_AC_only),
                ("V7 (V1 ∩ V4)", m_AD_only),
                ("V8 (V2 ∩ V3)", m_BC_only),
                ("V9 (V2 ∩ V4)", m_BD_only),
                ("V10 (V3 ∩ V4)", m_CD_only),
                ("V11 (V1 ∩ V2 ∩ V3)", m_ABC_only),
                ("V12 (V1 ∩ V2 ∩ V4)", m_ABD_only),
                ("V13 (V1 ∩ V3 ∩ V4)", m_ACD_only),
                ("V14 (V2 ∩ V3 ∩ V4)", m_BCD_only),
                ("V15 (BULLSEYE: All Four)", m_ABCD),
            ]

            for i, (label, mask) in enumerate(bidang_data):
                col = export_grids[i % 4]
                sub_df = df_v[mask].copy()
                cnt = len(sub_df)
                with col:
                    if cnt > 0:
                        excel_data = get_raw_excel(sub_df)
                        st.download_button(
                            label=f"📥 {label} ({cnt})",
                            data=excel_data,
                            file_name=f"Venn_{label.split(' ')[0]}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                    else:
                        st.button(f"📥 {label} (0)", disabled=True, use_container_width=True, key=f"btn_venn_{i}")

            st.markdown("---")
            st.markdown("### 📋 BULLSEYE Table Preview (V15)")

            df_bullseye = df_v[m_ABCD].copy()
            if df_bullseye.empty:
                st.warning("No SKUs met the BULLSEYE criteria (V15 Intersection).")
            else:
                base_display_cols = ['SKU', 'Product Name', 'New Code', 'New Product Name', 'Remark', 'Remark2',
                                     'Status', 'Country', 'Source_Sheet', 'Gross Sales (Current)', 'Qty (Current)',
                                     'Qty Growth (%)', 'Gross Margin (Current)', 'Gross Margin (%)',
                                     'Contribution Margin (Current)', 'Contribution Margin (%)']
                display_cols = [c for c in base_display_cols if c in df_bullseye.columns]
                df_export_bullseye = df_bullseye[display_cols].copy()

                format_dict_venn = {
                    'Gross Sales (Current)': 'Rp {:,.0f}',
                    'Qty (Current)': '{:,.0f}',
                    'Qty Growth (%)': '{:.2f}%',
                    'Gross Margin (Current)': 'Rp {:,.0f}',
                    'Gross Margin (%)': '{:.2f}%',
                    'Contribution Margin (Current)': 'Rp {:,.0f}',
                    'Contribution Margin (%)': '{:.2f}%'
                }
                st.dataframe(df_export_bullseye.style.format(format_dict_venn, na_rep=""), use_container_width=True,
                             hide_index=True)

        # ---------------------------------------------------------
        # TAB 3: RATIONALIZATION PROGRESS MONITOR
        # ---------------------------------------------------------
        with tab_progress:
            st.markdown("### 📉 SKU Rationalization Progress Monitor")

            # --- DYNAMIC CALCULATION FROM RAW DATA (df_raw) ---
            df_uncommon = df_raw.copy()

            def get_sum(df, col):
                return df[col].sum() if col in df.columns else 0

            # Phase 1: Initial Base (All data without exception, Uncommon)
            p1_sku = len(df_uncommon)
            p1_sales = get_sum(df_uncommon, 'Gross Sales (Current)')
            p1_gm = get_sum(df_uncommon, 'Gross Margin (Current)')
            p1_cm = get_sum(df_uncommon, 'Contribution Margin (Current)')

            # Phase 2: House Keeping In-Active 768 SKU
            if 'Remark' in df_uncommon.columns:
                remark_clean = df_uncommon['Remark'].astype(str).str.strip().str.upper()
                mask_exc = remark_clean.isin(['DISC', 'RENEWAL'])
                df_p2 = df_uncommon[~mask_exc]
            else:
                df_p2 = df_uncommon.copy()

            p2_sku = len(df_p2)
            p2_sales = get_sum(df_p2, 'Gross Sales (Current)')
            p2_gm = get_sum(df_p2, 'Gross Margin (Current)')
            p2_cm = get_sum(df_p2, 'Contribution Margin (Current)')

            # Phase 3: Discontinue 50 SKU
            if 'Status' in df_p2.columns:
                mask_deact = df_p2['Status'].astype(str).str.contains('Deactivated 10-Aug', case=False, na=False)
                df_p3 = df_p2[~mask_deact]
            else:
                df_p3 = df_p2.copy()

            p3_sku = len(df_p3)
            p3_sales = get_sum(df_p3, 'Gross Sales (Current)')
            p3_gm = get_sum(df_p3, 'Gross Margin (Current)')
            p3_cm = get_sum(df_p3, 'Contribution Margin (Current)')

            # Phase 4: Commonized SKU Export from 213 to 81 SKU
            df_p4_raw = df_uncommon.copy()
            if 'Status' in df_p4_raw.columns:
                df_p4_raw = df_p4_raw[df_p4_raw['Status'].astype(str).str.strip().str.upper() == 'ACTIVE']

            if 'New Code' in df_p4_raw.columns and 'New Product Name' in df_p4_raw.columns:
                df_p4_raw['New Code'] = df_p4_raw['New Code'].fillna('UNKNOWN').astype(str)
                df_p4_raw['New Product Name'] = df_p4_raw['New Product Name'].fillna('UNKNOWN').astype(str)

                group_cols_p4 = ['Source_Sheet', 'New Code', 'New Product Name']
                num_cols_p4 = df_p4_raw.select_dtypes(include=[np.number]).columns.tolist()
                sum_cols_p4 = [c for c in num_cols_p4 if not c.endswith('(%)')]

                agg_dict_p4 = {c: 'sum' for c in sum_cols_p4}
                df_p4_common = df_p4_raw.groupby(group_cols_p4, as_index=False).agg(agg_dict_p4)
            else:
                df_p4_common = df_p4_raw.copy()

            p4_sku = len(df_p4_common)
            p4_sales = get_sum(df_p4_common, 'Gross Sales (Current)')
            p4_gm = get_sum(df_p4_common, 'Gross Margin (Current)')
            p4_cm = get_sum(df_p4_common, 'Contribution Margin (Current)')

            # Phase 5: 90% up Gross Sales, GM and CM
            if not df_p4_common.empty and 'Gross Sales (Current)' in df_p4_common.columns:
                cutoff = max(1, round(len(df_p4_common) * 0.45))
                df_p5 = df_p4_common.sort_values(by='Gross Sales (Current)', ascending=False).head(cutoff)
            else:
                df_p5 = df_p4_common.copy()

            p5_sku = len(df_p5)
            p5_sales = get_sum(df_p5, 'Gross Sales (Current)')
            p5_gm = get_sum(df_p5, 'Gross Margin (Current)')
            p5_cm = get_sum(df_p5, 'Contribution Margin (Current)')

            # Render chart data
            progress_data_numeric = pd.DataFrame({
                'Phase': [
                    'P1 (Initial Base)',
                    'P2 (Housekeeping In-Active 768 SKUs)',
                    'P3 (Discontinue 50 SKUs)',
                    'P4 (Commonized SKU Export from 213 to 81 SKUs)',
                    'P5 (Top Performing SKUs - 90% Cumulative or 45% Count)'
                ],
                'Jumlah SKU': [p1_sku, p2_sku, p3_sku, p4_sku, p5_sku],
                'Gross Sales (IDR)': [p1_sales, p2_sales, p3_sales, p4_sales, p5_sales],
                'Gross Margin (IDR)': [p1_gm, p2_gm, p3_gm, p4_gm, p5_gm],
                'Contribution Margin (IDR)': [p1_cm, p2_cm, p3_cm, p4_cm, p5_cm]
            })

            # Helper for formatting table using comma as thousands separator
            def format_id_rupiah(val):
                return f"{int(val):,}"

            display_progress_df = progress_data_numeric.copy()
            display_progress_df.rename(columns={'Jumlah SKU': 'Total SKUs'}, inplace=True)
            display_progress_df['Total SKUs'] = display_progress_df['Total SKUs'].apply(format_id_rupiah)
            display_progress_df['Gross Sales (IDR)'] = display_progress_df['Gross Sales (IDR)'].apply(format_id_rupiah)
            display_progress_df['Gross Margin (IDR)'] = display_progress_df['Gross Margin (IDR)'].apply(
                format_id_rupiah)
            display_progress_df['Contribution Margin (IDR)'] = display_progress_df['Contribution Margin (IDR)'].apply(
                format_id_rupiah)

            st.dataframe(display_progress_df, use_container_width=True, hide_index=True)

            progress_metric = st.selectbox(
                "Select Tracking Metric for Y-Axis:",
                ["Total SKUs", "Gross Sales (IDR)", "Gross Margin (IDR)", "Contribution Margin (IDR)"]
            )

            # Map the metric back to the dataframe column for plotting
            plot_metric_col = progress_metric if progress_metric != "Total SKUs" else "Jumlah SKU"

            # Executive Colors
            executive_colors = ['#1e3a8a', '#d97706', '#0f766e', '#6b21a8', '#be123c']

            # Rendering the Chart
            fig_prog = px.bar(
                progress_data_numeric,
                x='Phase',
                y=plot_metric_col,
                color='Phase',
                color_discrete_sequence=executive_colors,
                title=f"<b>RATIONALIZATION LIFECYCLE: {progress_metric.upper()}</b>"
            )

            if plot_metric_col == "Jumlah SKU":
                fig_prog.update_traces(
                    texttemplate='<b>%{y:,.0f} SKUs</b>',
                    textposition='outside',
                    textfont=dict(size=16, color='#0f172a'),
                    width=0.55, marker_line_width=0
                )
            else:
                fig_prog.update_traces(
                    texttemplate='<b>Rp %{y:,.0f}</b>',
                    textposition='outside',
                    textfont=dict(size=16, color='#0f172a'),
                    width=0.55, marker_line_width=0
                )

            fig_prog.update_layout(
                separators=",.",  # Plotly uses dot (.) as thousands separator and comma (,) as decimal
                plot_bgcolor='white',
                paper_bgcolor='white',
                showlegend=False,
                yaxis=dict(
                    showgrid=True, gridcolor='#f1f5f9',
                    showline=False, showticklabels=False, title=""
                ),
                xaxis=dict(
                    showgrid=False, showline=True, linewidth=2, linecolor='#334155',
                    title="", tickfont=dict(size=13, color='#334155', family="Arial")
                ),
                title_font=dict(size=22, color='#0f172a', family="Arial"),
                margin=dict(t=80, b=40, l=40, r=40),
                height=500
            )

            # Ensure Y-axis range is high enough so outside text doesn't get clipped
            max_y = progress_data_numeric[plot_metric_col].max()
            fig_prog.update_yaxes(range=[0, max_y * 1.25])

            st.plotly_chart(fig_prog, use_container_width=True)

            # =========================================================
            # NEW MODULE: FINANCIAL OUTLOOK (P&L MONITORING)
            # =========================================================
            st.markdown("---")
            st.markdown("### 📋 Financial Outlook: P&L Monitoring (Surviving SKUs)")
            st.info(
                "💡 Upload **3 separate CSV files**. The system uses the **Base P&L Data** to establish the initial benchmark (full historical granularity), and applies the **SKU Mapping File** to the **Quarterly P&L Data** to project the consolidated Phase 5 performance.")

            col_up1, col_up2, col_up3 = st.columns(3)
            with col_up1:
                pl_base_file = st.file_uploader("1. Base P&L Data (.csv)", type=["csv"], key="pl_base_up")
            with col_up2:
                pl_proj_file = st.file_uploader("2. Quarterly P&L Data (.csv)", type=["csv"], key="pl_proj_up")
            with col_up3:
                pl_map_file = st.file_uploader("3. SKU Mapping (.csv)", type=["csv"], key="pl_map_up")

            target_gm_pct = st.number_input(
                "Target Gross Margin Optimization (%)",
                min_value=0.0,
                max_value=100.0,
                value=0.0,
                step=0.1,
                help="Set to 0.0% for As-Is Scenario (no price changes). Set a value > 0% to optimize prices for underperforming SKUs."
            )

            if pl_base_file is not None and pl_proj_file is not None and pl_map_file is not None:
                try:
                    # 0. READ MAPPING FILE AND CREATE DICTIONARIES
                    df_map = pd.read_csv(pl_map_file, sep=None, engine='python', dtype=str)
                    df_map.columns = df_map.columns.str.strip()
                    for c in df_map.columns:
                        df_map[c] = df_map[c].astype(str).str.strip().replace(['nan', 'NaN', 'None'], '')

                    # Target correct columns using Product Key to avoid whitespace typo issues
                    col_pk = next((col for col in df_map.columns if 'Produk_Key' in col or 'Product_Key' in col),
                                  df_map.columns[0])
                    col_nc = next((col for col in df_map.columns if 'New Code' in col), df_map.columns[2])
                    col_npn = next((col for col in df_map.columns if 'New Product Name' in col), df_map.columns[3])

                    dict_code = dict(zip(df_map[col_pk], df_map[col_nc]))
                    dict_name = dict(zip(df_map[col_pk], df_map[col_npn]))

                    # HELPER FUNCTION TO LOAD AND CLEAN P&L DATA
                    def load_and_clean_pl(file, apply_mapping=False):
                        # 1. READ FILE AS TEXT FOR IDENTITY COLUMNS TO PREVENT SCIENTIFIC NOTATION
                        df = pd.read_csv(file, sep=';', dtype=str)
                        if len(df.columns) == 1:
                            file.seek(0)
                            df = pd.read_csv(file, sep=',', dtype=str)

                        df.columns = df.columns.str.strip()

                        # 2. TEXT DATA CLEANSING
                        text_cols_pl = ['Produk_Key', 'Product Name', 'Brand', 'Series', 'Segmentation', 'Category',
                                        'Market', 'Month', 'Year']
                        for c in text_cols_pl:
                            if c in df.columns:
                                df[c] = df[c].astype(str).str.replace(r'\.0$', '', regex=True).str.strip().replace(
                                    ['nan', 'NaN', 'None'], '')

                        # 3. NUMERIC DATA CLEANSING (Handling Currency Formats)
                        def clean_currency(x):
                            if pd.isna(x): return 0.0
                            x_str = str(x).strip()
                            if x_str in ['-', '']: return 0.0
                            x_str = x_str.replace('.', '').replace(',', '.')
                            try:
                                return float(x_str)
                            except:
                                return 0.0

                        cols_to_clean = [
                            'Gross Sales', 'Sales Return', 'Sales Deduction', 'Net Sales', 'Total COGS', 'COGS_Regular',
                            'Royalty', 'Gross Profit',
                            'Advertising Activity', 'Domestic Marketing', 'Export Marketing Support',
                            'Export Marketing',
                            'Research & Analysis Expenses', 'Selling Activity Expenses', 'Travelling Expenses',
                            'Others Expenses',
                            'Selling Personnel', 'Export Selling Activity', 'Export Transportation', 'Export Personnel',
                            'Export Others', 'General Affair Expenses', 'Finance Expenses', 'Computer Expenses',
                            'General Personnel Expenses', 'Education', 'Entertainment', 'Sundry Expenses',
                            'General Administrative', 'Depreciation/Amortisation'
                        ]
                        for c in cols_to_clean:
                            if c in df.columns:
                                df[c] = df[c].apply(clean_currency)

                        # INJECT MAPPING INTO P&L DATA VIA ABSOLUTE PRODUCT KEY (ONLY IF REQUESTED)
                        if apply_mapping:
                            if 'Produk_Key' in df.columns:
                                df['New Code'] = df['Produk_Key'].map(dict_code).fillna(df.get('Produk_Key', ''))
                                df['New Product Name'] = df['Produk_Key'].map(dict_name).fillna(
                                    df.get('Product Name', ''))
                            elif 'Product Name' in df.columns:
                                # Fallback if Produk_Key is absent
                                df['New Code'] = df['Product Name'].map(dict_code).fillna(df.get('Produk_Key', ''))
                                df['New Product Name'] = df['Product Name'].map(dict_name).fillna(df['Product Name'])

                        return df

                    # Load Base without mapping (to retain 2424 SKUs natively), Load Projected with mapping
                    df_pl_base = load_and_clean_pl(pl_base_file, apply_mapping=False)
                    df_pl_proj = load_and_clean_pl(pl_proj_file, apply_mapping=True)

                    # 4. CHRONOLOGICAL QUARTER MAPPING FOR PROJECTED DATA
                    if 'Month' in df_pl_proj.columns and 'Year' in df_pl_proj.columns:
                        df_pl_proj['Month_Str'] = df_pl_proj['Month'].astype(str).str.strip().str.capitalize()
                        df_pl_proj['Year_Num'] = pd.to_numeric(df_pl_proj['Year'], errors='coerce').fillna(2000)

                        df_pl_proj['Month_Num'] = pd.to_datetime(df_pl_proj['Month_Str'], format='%B',
                                                                 errors='coerce').dt.month.fillna(1)
                        unique_periods = df_pl_proj[
                            ['Year_Num', 'Month_Num', 'Month_Str']].drop_duplicates().sort_values(
                            by=['Year_Num', 'Month_Num'])

                        chunks = np.array_split(unique_periods['Month_Str'].values, 4)
                        month_to_q = {}
                        for i, chunk in enumerate(chunks):
                            for m in chunk:
                                month_to_q[m] = f"Q{i + 1}"

                        df_pl_proj['Quarter'] = df_pl_proj['Month_Str'].map(month_to_q).fillna('Q1')
                    else:
                        df_pl_proj['Quarter'] = 'Q1'

                    # 5. ISOLATE THE DEFINITIVE 663 MASTER SKUs
                    # Obtain exactly one unique row per Master SKU name from P5
                    surviving_masters_df = df_p5[['New Code', 'New Product Name', 'Source_Sheet']].drop_duplicates(
                        subset=['New Product Name'])
                    surviving_masters_df.rename(columns={'Source_Sheet': 'Market'}, inplace=True)
                    surviving_masters_list = surviving_masters_df['New Product Name'].tolist()

                    df_pl_survivors = df_pl_proj[df_pl_proj['New Product Name'].isin(surviving_masters_list)].copy()

                    # 6. SCENARIO PROJECTION LOGIC
                    df_pl_survivors['Calc_GS_Net'] = df_pl_survivors.get('Gross Sales', 0) + df_pl_survivors.get(
                        'Sales Return', 0)

                    # Grouping strictly by New Product Name to lift prices uniformly at the Master SKU level
                    sku_gm = df_pl_survivors.groupby('New Product Name').agg(
                        total_gs=('Calc_GS_Net', 'sum'),
                        total_cogs_reg=('COGS_Regular', 'sum'),
                        total_royalty=('Royalty', 'sum')
                    )
                    sku_gm['total_gm'] = sku_gm['total_gs'] - sku_gm['total_cogs_reg'] - sku_gm['total_royalty']
                    sku_gm['gm_pct'] = np.where(sku_gm['total_gs'] > 0, sku_gm['total_gm'] / sku_gm['total_gs'], 0)

                    df_proj = df_pl_survivors.copy()

                    if target_gm_pct > 0.0:
                        target_gm_decimal = target_gm_pct / 100.0
                        skus_to_improve = sku_gm[sku_gm['gm_pct'] < target_gm_decimal].index
                        target_gs = (sku_gm.loc[skus_to_improve, 'total_cogs_reg'] + sku_gm.loc[
                            skus_to_improve, 'total_royalty']) / (1 - target_gm_decimal)

                        # 100% Multiplier for Gross Sales
                        gs_mult = target_gs / (sku_gm.loc[skus_to_improve, 'total_gs'] + 1e-9)
                        df_proj['gs_mult'] = df_proj['New Product Name'].map(gs_mult).fillna(1.0)
                    else:
                        df_proj['gs_mult'] = 1.0

                    # Apply multiplier to transaction rows (Only Gross Sales increases)
                    row_gs_delta = df_proj.get('Gross Sales', 0) * (df_proj['gs_mult'] - 1)
                    if 'Gross Sales' in df_proj.columns: df_proj['Gross Sales'] += row_gs_delta
                    if 'Net Sales' in df_proj.columns: df_proj['Net Sales'] += row_gs_delta

                    # UPDATE GROSS PROFIT ACCORDINGLY
                    if 'Gross Profit' in df_proj.columns:
                        df_proj['Gross Profit'] += row_gs_delta

                    # 7. CORE CALCULATION ENGINE
                    sga_cols = [
                        'Selling Activity Expenses', 'Travelling Expenses', 'Others Expenses', 'Selling Personnel',
                        'Export Selling Activity', 'Export Marketing Support', 'Export Transportation',
                        'Export Personnel',
                        'Export Others',
                        'General Affair Expenses', 'Finance Expenses', 'Computer Expenses',
                        'General Personnel Expenses',
                        'Education', 'Entertainment', 'Sundry Expenses', 'Depreciation/Amortisation'
                    ]

                    def calculate_metrics(df_source):
                        res = {}
                        res['Calc_Gross_Sales'] = df_source.get('Gross Sales', 0).sum()
                        res['Calc_Sales_Ded'] = df_source.get('Sales Deduction', 0).sum() + df_source.get(
                            'Sales Return', 0).sum()
                        res['Calc_Net_Sales'] = df_source.get('Net Sales', 0).sum()
                        res['Calc_COGS'] = df_source.get('Total COGS', 0).sum()
                        res['Calc_Gross_Profit'] = df_source.get('Gross Profit', 0).sum()

                        res['Calc_AP'] = abs(df_source.get('Advertising Activity', 0).sum())
                        res['Calc_RD'] = abs(df_source.get('Research & Analysis Expenses', 0).sum())

                        sga_sum = 0
                        for c in sga_cols:
                            if c in df_source.columns:
                                sga_sum += df_source[c].sum()
                        res['Calc_SGA'] = abs(sga_sum)
                        res['Calc_DA'] = abs(df_source.get('Depreciation/Amortisation', 0).sum())
                        return pd.Series(res)

                    def get_gm_ratio(df_source):
                        gs = df_source.get('Gross Sales', 0).sum() + df_source.get('Sales Return', 0).sum()
                        gm = gs - df_source.get('COGS_Regular', 0).sum() - df_source.get('Royalty', 0).sum()
                        if gs == 0: return 0.0
                        return gm / gs

                    # 8. AGGREGATION & FINAL TABLE PREPARATION
                    base_metrics = calculate_metrics(df_pl_base)

                    quarters = ['Q1', 'Q2', 'Q3', 'Q4']
                    q_metrics = {}
                    for q in quarters:
                        q_metrics[q] = calculate_metrics(df_proj[df_proj['Quarter'] == q])

                    pl_compiled = pd.DataFrame(q_metrics)
                    for q in quarters:
                        if q not in pl_compiled.columns:
                            pl_compiled[q] = 0.0

                    pl_compiled['Base'] = base_metrics
                    pl_compiled['FY'] = pl_compiled[['Q1', 'Q2', 'Q3', 'Q4']].sum(axis=1)

                    pl_bn = pl_compiled / 1e9

                    def get_row(metric_name):
                        return pl_bn.loc[metric_name] if metric_name in pl_bn.index else pd.Series([0] * 6,
                                                                                                   index=['Base', 'Q1',
                                                                                                          'Q2', 'Q3',
                                                                                                          'Q4', 'FY'])

                    r_gs = get_row('Calc_Gross_Sales')
                    r_sd = get_row('Calc_Sales_Ded')
                    r_ns = get_row('Calc_Net_Sales')
                    r_cogs = get_row('Calc_COGS')
                    r_gp = get_row('Calc_Gross_Profit')
                    r_ap = get_row('Calc_AP')
                    r_rd = get_row('Calc_RD')
                    r_sga = get_row('Calc_SGA')
                    r_da = get_row('Calc_DA')

                    r_ebit = r_gp - r_ap - r_rd - r_sga
                    r_ebitda = r_ebit + r_da

                    r_gm_ratio = {
                        'Base': get_gm_ratio(df_pl_base),
                        'Q1': get_gm_ratio(df_proj[df_proj['Quarter'] == 'Q1']),
                        'Q2': get_gm_ratio(df_proj[df_proj['Quarter'] == 'Q2']),
                        'Q3': get_gm_ratio(df_proj[df_proj['Quarter'] == 'Q3']),
                        'Q4': get_gm_ratio(df_proj[df_proj['Quarter'] == 'Q4']),
                        'FY': get_gm_ratio(df_proj[df_proj['Quarter'].isin(['Q1', 'Q2', 'Q3', 'Q4'])])
                    }

                    # --- 9. SKU DETAILS DATAFRAME CREATION ---

                    # 9A. BASE details use the unmapped, raw original data
                    def get_sku_details_base(df_source):
                        if df_source.empty: return pd.DataFrame()

                        # Group by original fields if they exist
                        group_keys = [c for c in ['Produk_Key', 'Product Name', 'Market'] if c in df_source.columns]
                        if not group_keys: return pd.DataFrame()

                        def calc_row(x):
                            gs_pure = x.get('Gross Sales', 0).sum()
                            sales_ded_tot = x.get('Sales Deduction', 0).sum() + x.get('Sales Return', 0).sum()
                            ns = x.get('Net Sales', 0).sum()
                            total_cogs = x.get('Total COGS', 0).sum()
                            gp_pl = x.get('Gross Profit', 0).sum()

                            cogs_reg = x.get('COGS_Regular', 0).sum()
                            royalty = x.get('Royalty', 0).sum()
                            gs_net = x.get('Gross Sales', 0).sum() + x.get('Sales Return', 0).sum()
                            gm_custom = gs_net - cogs_reg - royalty

                            ap = abs(x.get('Advertising Activity', 0).sum())
                            rd = abs(x.get('Research & Analysis Expenses', 0).sum())
                            sga = abs(sum(x.get(c, 0).sum() for c in sga_cols if c in x.columns))
                            da = abs(x.get('Depreciation/Amortisation', 0).sum())

                            return pd.Series({
                                'Gross Sales': gs_pure,
                                'Sales Deductions': sales_ded_tot,
                                'Net Sales': ns,
                                '(-) Cost of Sales': total_cogs,
                                'Gross Profit': gp_pl,
                                '(-) A&P': ap,
                                '(-) R&D': rd,
                                '(-) Other SG&A': sga,
                                'EBIT': gp_pl - ap - rd - sga,
                                '(+) D&A': da,
                                'EBITDA': gp_pl - ap - rd - sga + da,
                                'Ratio Gross Margin': (gm_custom / gs_net) if gs_net != 0 else 0
                            })

                        res = df_source.groupby(group_keys).apply(calc_row).reset_index()
                        res.insert(0, 'Period', 'Base')

                        # Standardize columns to merge cleanly with projected data
                        if 'Produk_Key' in res.columns:
                            res.rename(columns={'Produk_Key': 'SKU'}, inplace=True)
                        return res

                    # 9B. PROJECTED details map directly to exactly 663 Master SKUs
                    def get_sku_details_proj(df_source, period_name):
                        def calc_row(x):
                            gs_pure = x.get('Gross Sales', 0).sum()
                            sales_ded_tot = x.get('Sales Deduction', 0).sum() + x.get('Sales Return', 0).sum()
                            ns = x.get('Net Sales', 0).sum()
                            total_cogs = x.get('Total COGS', 0).sum()
                            gp_pl = x.get('Gross Profit', 0).sum()

                            cogs_reg = x.get('COGS_Regular', 0).sum()
                            royalty = x.get('Royalty', 0).sum()
                            gs_net = x.get('Gross Sales', 0).sum() + x.get('Sales Return', 0).sum()
                            gm_custom = gs_net - cogs_reg - royalty

                            ap = abs(x.get('Advertising Activity', 0).sum())
                            rd = abs(x.get('Research & Analysis Expenses', 0).sum())
                            sga = abs(sum(x.get(c, 0).sum() for c in sga_cols if c in x.columns))
                            da = abs(x.get('Depreciation/Amortisation', 0).sum())

                            return pd.Series({
                                'Gross Sales': gs_pure,
                                'Sales Deductions': sales_ded_tot,
                                'Net Sales': ns,
                                '(-) Cost of Sales': total_cogs,
                                'Gross Profit': gp_pl,
                                '(-) A&P': ap,
                                '(-) R&D': rd,
                                '(-) Other SG&A': sga,
                                'EBIT': gp_pl - ap - rd - sga,
                                '(+) D&A': da,
                                'EBITDA': gp_pl - ap - rd - sga + da,
                                'Ratio Gross Margin': (gm_custom / gs_net) if gs_net != 0 else 0
                            })

                        if df_source.empty:
                            res_final = surviving_masters_df.copy()
                            for col in ['Gross Sales', 'Sales Deductions', 'Net Sales', '(-) Cost of Sales',
                                        'Gross Profit', '(-) A&P', '(-) R&D', '(-) Other SG&A', 'EBIT', '(+) D&A',
                                        'EBITDA', 'Ratio Gross Margin']:
                                res_final[col] = 0.0
                            res_final.insert(0, 'Period', period_name)
                            res_final.rename(columns={'New Code': 'SKU', 'New Product Name': 'Product Name'},
                                             inplace=True)
                            return res_final

                        # Group strictly by New Product Name to avoid mismatch due to New Code variants
                        res = df_source.groupby('New Product Name').apply(calc_row).reset_index()

                        # Left join enforces EXACTLY the 663 Master SKUs to appear, substituting 0.0 for missing ones
                        res_final = surviving_masters_df.merge(res, on='New Product Name', how='left').fillna(0.0)
                        res_final.insert(0, 'Period', period_name)
                        res_final.rename(columns={'New Code': 'SKU', 'New Product Name': 'Product Name'}, inplace=True)
                        return res_final

                    details_base = get_sku_details_base(df_pl_base)
                    details_q1 = get_sku_details_proj(df_proj[df_proj['Quarter'] == 'Q1'], 'Q1')
                    details_q2 = get_sku_details_proj(df_proj[df_proj['Quarter'] == 'Q2'], 'Q2')
                    details_q3 = get_sku_details_proj(df_proj[df_proj['Quarter'] == 'Q3'], 'Q3')
                    details_q4 = get_sku_details_proj(df_proj[df_proj['Quarter'] == 'Q4'], 'Q4')
                    details_fy = get_sku_details_proj(df_proj, 'FY')

                    df_all_details = pd.concat(
                        [details_base, details_q1, details_q2, details_q3, details_q4, details_fy], ignore_index=True)

                    # 10. HTML INJECTION FOR CORPORATE TABLE AESTHETICS
                    def fmt(val):
                        if val == 0: return "-"
                        return f"{val:,.1f}"

                    def fmt_pct(val):
                        if val == 0: return "0.0%"
                        return f"{val * 100:.1f}%"

                    st.caption(f"*(All figures are presented in Billion Rupiah / IDR Bn)*")

                    table_title = "Financial Model & Tracking (As-Is Scenario)"
                    if target_gm_pct > 0.0:
                        table_title = f"Financial Model & Tracking (Optimized Target GM {target_gm_pct:.1f}%)"

                    html_table = f"""
                    <style>
                    .pl-table {{ width: 100%; border-collapse: collapse; font-family: 'Segoe UI', sans-serif; font-size: 14px; background-color: white; margin-bottom: 5px; }}
                    .pl-table th {{ background-color: #4472c4; color: white; padding: 8px; text-align: center; border: 1px solid #b4c6e7; font-weight: normal; }}
                    .pl-table th:first-child {{ text-align: left; }}
                    .pl-table td {{ padding: 6px 8px; border: 1px solid #b4c6e7; text-align: right; color: #1f2937; }}
                    .pl-table td:first-child {{ text-align: left; font-weight: normal; }}
                    .pl-bold-row td {{ font-weight: bold !important; border-top: 2px solid #4472c4; border-bottom: 2px solid #4472c4; }}
                    .pl-ratio-row td {{ font-weight: bold !important; background-color: #f8fafc; font-style: italic; border-top: 2px solid #94a3b8; border-bottom: 2px solid #94a3b8; }}
                    </style>
                    <table class='pl-table'>
                        <tr><th colspan="7" style="font-size: 16px; font-weight: bold; background-color: #4472c4; color: white; text-align: center;">{table_title}</th></tr>
                        <tr><th>IDR Bn</th><th>Base</th><th>Q1</th><th>Q2</th><th>Q3</th><th>Q4</th><th>FY</th></tr>
                        <tr><td>Gross Sales</td><td>{fmt(r_gs['Base'])}</td><td>{fmt(r_gs['Q1'])}</td><td>{fmt(r_gs['Q2'])}</td><td>{fmt(r_gs['Q3'])}</td><td>{fmt(r_gs['Q4'])}</td><td>{fmt(r_gs['FY'])}</td></tr>
                        <tr><td>Sales Deductions</td><td>{fmt(r_sd['Base'])}</td><td>{fmt(r_sd['Q1'])}</td><td>{fmt(r_sd['Q2'])}</td><td>{fmt(r_sd['Q3'])}</td><td>{fmt(r_sd['Q4'])}</td><td>{fmt(r_sd['FY'])}</td></tr>
                        <tr class='pl-bold-row'><td>Net Sales</td><td>{fmt(r_ns['Base'])}</td><td>{fmt(r_ns['Q1'])}</td><td>{fmt(r_ns['Q2'])}</td><td>{fmt(r_ns['Q3'])}</td><td>{fmt(r_ns['Q4'])}</td><td>{fmt(r_ns['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr><td>(-) Cost of Sales</td><td>{fmt(r_cogs['Base'])}</td><td>{fmt(r_cogs['Q1'])}</td><td>{fmt(r_cogs['Q2'])}</td><td>{fmt(r_cogs['Q3'])}</td><td>{fmt(r_cogs['Q4'])}</td><td>{fmt(r_cogs['FY'])}</td></tr>
                        <tr class='pl-bold-row'><td>Gross Profit</td><td>{fmt(r_gp['Base'])}</td><td>{fmt(r_gp['Q1'])}</td><td>{fmt(r_gp['Q2'])}</td><td>{fmt(r_gp['Q3'])}</td><td>{fmt(r_gp['Q4'])}</td><td>{fmt(r_gp['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr><td>(-) A&P</td><td>{fmt(r_ap['Base'])}</td><td>{fmt(r_ap['Q1'])}</td><td>{fmt(r_ap['Q2'])}</td><td>{fmt(r_ap['Q3'])}</td><td>{fmt(r_ap['Q4'])}</td><td>{fmt(r_ap['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr><td>(-) R&D</td><td>{fmt(r_rd['Base'])}</td><td>{fmt(r_rd['Q1'])}</td><td>{fmt(r_rd['Q2'])}</td><td>{fmt(r_rd['Q3'])}</td><td>{fmt(r_rd['Q4'])}</td><td>{fmt(r_rd['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr><td>(-) Other SG&A</td><td>{fmt(r_sga['Base'])}</td><td>{fmt(r_sga['Q1'])}</td><td>{fmt(r_sga['Q2'])}</td><td>{fmt(r_sga['Q3'])}</td><td>{fmt(r_sga['Q4'])}</td><td>{fmt(r_sga['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr class='pl-bold-row'><td>EBIT</td><td>{fmt(r_ebit['Base'])}</td><td>{fmt(r_ebit['Q1'])}</td><td>{fmt(r_ebit['Q2'])}</td><td>{fmt(r_ebit['Q3'])}</td><td>{fmt(r_ebit['Q4'])}</td><td>{fmt(r_ebit['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr><td>(+) D&A</td><td>{fmt(r_da['Base'])}</td><td>{fmt(r_da['Q1'])}</td><td>{fmt(r_da['Q2'])}</td><td>{fmt(r_da['Q3'])}</td><td>{fmt(r_da['Q4'])}</td><td>{fmt(r_da['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr class='pl-bold-row'><td>EBITDA</td><td>{fmt(r_ebitda['Base'])}</td><td>{fmt(r_ebitda['Q1'])}</td><td>{fmt(r_ebitda['Q2'])}</td><td>{fmt(r_ebitda['Q3'])}</td><td>{fmt(r_ebitda['Q4'])}</td><td>{fmt(r_ebitda['FY'])}</td></tr>
                        <tr><td colspan='7'></td></tr>
                        <tr class='pl-ratio-row'><td>Ratio Gross Margin</td><td>{fmt_pct(r_gm_ratio['Base'])}</td><td>{fmt_pct(r_gm_ratio['Q1'])}</td><td>{fmt_pct(r_gm_ratio['Q2'])}</td><td>{fmt_pct(r_gm_ratio['Q3'])}</td><td>{fmt_pct(r_gm_ratio['Q4'])}</td><td>{fmt_pct(r_gm_ratio['FY'])}</td></tr>
                    </table>
                    """
                    st.markdown(html_table, unsafe_allow_html=True)

                    if target_gm_pct > 0.0:
                        st.caption(
                            f"*The table above compares **Base** (displaying all raw historical SKUs) with **Projected FY** (total aggregation of strictly {len(surviving_masters_list)} Master SKUs configured via the mapping, targeted to a minimum {target_gm_pct:.1f}% Margin purely through price adjustments).*")
                    else:
                        st.caption(
                            f"*The table above compares **Base** (displaying all raw historical SKUs) with **Projected FY** (total aggregation of strictly {len(surviving_masters_list)} Master SKUs configured via the mapping, reflecting actual historical performance without any price adjustments).*")

                    # 11. EXPORT TO EXCEL WITH IDENTICAL STYLING & EXACT SKU DETAILS
                    export_data = [
                        ['Gross Sales', r_gs['Base'], r_gs['Q1'], r_gs['Q2'], r_gs['Q3'], r_gs['Q4'], r_gs['FY']],
                        ['Sales Deductions', r_sd['Base'], r_sd['Q1'], r_sd['Q2'], r_sd['Q3'], r_sd['Q4'], r_sd['FY']],
                        ['Net Sales', r_ns['Base'], r_ns['Q1'], r_ns['Q2'], r_ns['Q3'], r_ns['Q4'], r_ns['FY']],
                        ['', None, None, None, None, None, None],
                        ['(-) Cost of Sales', r_cogs['Base'], r_cogs['Q1'], r_cogs['Q2'], r_cogs['Q3'], r_cogs['Q4'],
                         r_cogs['FY']],
                        ['Gross Profit', r_gp['Base'], r_gp['Q1'], r_gp['Q2'], r_gp['Q3'], r_gp['Q4'], r_gp['FY']],
                        ['', None, None, None, None, None, None],
                        ['(-) A&P', r_ap['Base'], r_ap['Q1'], r_ap['Q2'], r_ap['Q3'], r_ap['Q4'], r_ap['FY']],
                        ['', None, None, None, None, None, None],
                        ['(-) R&D', r_rd['Base'], r_rd['Q1'], r_rd['Q2'], r_rd['Q3'], r_rd['Q4'], r_rd['FY']],
                        ['', None, None, None, None, None, None],
                        ['(-) Other SG&A', r_sga['Base'], r_sga['Q1'], r_sga['Q2'], r_sga['Q3'], r_sga['Q4'],
                         r_sga['FY']],
                        ['', None, None, None, None, None, None],
                        ['EBIT', r_ebit['Base'], r_ebit['Q1'], r_ebit['Q2'], r_ebit['Q3'], r_ebit['Q4'], r_ebit['FY']],
                        ['', None, None, None, None, None, None],
                        ['(+) D&A', r_da['Base'], r_da['Q1'], r_da['Q2'], r_da['Q3'], r_da['Q4'], r_da['FY']],
                        ['', None, None, None, None, None, None],
                        ['EBITDA', r_ebitda['Base'], r_ebitda['Q1'], r_ebitda['Q2'], r_ebitda['Q3'], r_ebitda['Q4'],
                         r_ebitda['FY']],
                        ['', None, None, None, None, None, None],
                        ['Ratio Gross Margin', r_gm_ratio['Base'], r_gm_ratio['Q1'], r_gm_ratio['Q2'],
                         r_gm_ratio['Q3'], r_gm_ratio['Q4'], r_gm_ratio['FY']]
                    ]

                    df_export_pl = pd.DataFrame(export_data, columns=['IDR Bn', 'Base', 'Q1', 'Q2', 'Q3', 'Q4', 'FY'])

                    buffer_pl = io.BytesIO()
                    with pd.ExcelWriter(buffer_pl, engine='xlsxwriter') as writer:
                        # WRITE SHEET 1: SUMMARY
                        df_export_pl.to_excel(writer, index=False, sheet_name='Financial_Outlook', startrow=1)
                        workbook = writer.book
                        worksheet = writer.sheets['Financial_Outlook']

                        header_format = workbook.add_format({
                            'bold': True, 'bg_color': '#4472c4', 'font_color': 'white',
                            'align': 'center', 'valign': 'vcenter', 'border': 1, 'border_color': '#b4c6e7'
                        })
                        header_left_format = workbook.add_format({
                            'bold': True, 'bg_color': '#4472c4', 'font_color': 'white',
                            'align': 'left', 'valign': 'vcenter', 'border': 1, 'border_color': '#b4c6e7'
                        })
                        title_format = workbook.add_format({
                            'bold': True, 'bg_color': '#4472c4', 'font_color': 'white',
                            'align': 'center', 'valign': 'vcenter', 'font_size': 14
                        })
                        data_format = workbook.add_format({
                            'num_format': '#,##0.1;-#,##0.1;"-"', 'border': 1, 'border_color': '#b4c6e7',
                            'valign': 'vcenter'
                        })
                        data_format_bold = workbook.add_format({
                            'bold': True, 'num_format': '#,##0.1;-#,##0.1;"-"', 'border': 1, 'border_color': '#b4c6e7',
                            'valign': 'vcenter', 'top': 2, 'bottom': 2, 'top_color': '#4472c4',
                            'bottom_color': '#4472c4'
                        })
                        label_format = workbook.add_format({
                            'border': 1, 'border_color': '#b4c6e7', 'valign': 'vcenter'
                        })
                        label_format_bold = workbook.add_format({
                            'bold': True, 'border': 1, 'border_color': '#b4c6e7', 'valign': 'vcenter',
                            'top': 2, 'bottom': 2, 'top_color': '#4472c4', 'bottom_color': '#4472c4'
                        })
                        ratio_format_bold = workbook.add_format({
                            'bold': True, 'bg_color': '#f8fafc', 'italic': True, 'border': 1, 'border_color': '#b4c6e7',
                            'valign': 'vcenter', 'top': 2, 'bottom': 2, 'top_color': '#94a3b8',
                            'bottom_color': '#94a3b8', 'align': 'right', 'num_format': '0.0%'
                        })
                        ratio_label_format = workbook.add_format({
                            'bold': True, 'bg_color': '#f8fafc', 'italic': True, 'border': 1, 'border_color': '#b4c6e7',
                            'valign': 'vcenter', 'top': 2, 'bottom': 2, 'top_color': '#94a3b8',
                            'bottom_color': '#94a3b8', 'align': 'left'
                        })

                        worksheet.merge_range('A1:G1', table_title, title_format)

                        for col_num, value in enumerate(df_export_pl.columns.values):
                            if col_num == 0:
                                worksheet.write(1, col_num, value, header_left_format)
                            else:
                                worksheet.write(1, col_num, value, header_format)

                        bold_rows = [2, 5, 13, 17]
                        ratio_row = 19

                        for row_num, row_data in enumerate(df_export_pl.values):
                            is_bold = row_num in bold_rows
                            is_ratio = (row_num == ratio_row)
                            is_empty = (row_data[0] == '')

                            for col_num, val in enumerate(row_data):
                                if pd.isna(val):
                                    val = ""
                                if col_num == 0:
                                    if is_ratio:
                                        worksheet.write(row_num + 2, col_num, val, ratio_label_format)
                                    else:
                                        worksheet.write(row_num + 2, col_num, val,
                                                        label_format_bold if is_bold else label_format)
                                else:
                                    if is_empty:
                                        worksheet.write(row_num + 2, col_num, "", label_format)
                                    elif is_ratio:
                                        worksheet.write(row_num + 2, col_num, val, ratio_format_bold)
                                    else:
                                        worksheet.write(row_num + 2, col_num, val,
                                                        data_format_bold if is_bold else data_format)

                        # WRITE SHEET 2: SKU DETAILS
                        df_all_details.to_excel(writer, index=False, sheet_name='SKU_Details')
                        ws_det = writer.sheets['SKU_Details']

                        num_fmt = workbook.add_format({'num_format': '#,##0', 'border': 1, 'border_color': '#d1d5db'})
                        pct_fmt = workbook.add_format(
                            {'num_format': '0.0%', 'border': 1, 'border_color': '#d1d5db', 'align': 'center',
                             'bold': True})
                        hdr_det_format = workbook.add_format(
                            {'bold': True, 'bg_color': '#4472c4', 'font_color': 'white', 'border': 1,
                             'align': 'center'})

                        # Dynamic column sizing
                        for col_idx, col_name in enumerate(df_all_details.columns):
                            if col_name in ['Period', 'Market']:
                                ws_det.set_column(col_idx, col_idx, 15)
                            elif col_name in ['SKU', 'Product Name']:
                                ws_det.set_column(col_idx, col_idx, 35)
                            elif col_name == 'Ratio Gross Margin':
                                ws_det.set_column(col_idx, col_idx, 18, pct_fmt)
                            else:
                                ws_det.set_column(col_idx, col_idx, 18, num_fmt)

                        for col_num, value in enumerate(df_all_details.columns.values):
                            ws_det.write(0, col_num, value, hdr_det_format)

                    col_dl_pl1, _ = st.columns([1, 4])
                    with col_dl_pl1:
                        st.download_button(
                            label="📥 Download P&L Table & Details (.xlsx)",
                            data=buffer_pl.getvalue(),
                            file_name="Financial_Outlook_Survivors.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            type="primary",
                            use_container_width=True
                        )

                except Exception as e:
                    st.error(
                        f"Failed to process P&L data. Ensure the table format meets standard requirements. Error details: {e}")


if __name__ == "__main__":
    main()
